#!/usr/bin/env python3
"""
Extract slide master background color
"""

from pptx import Presentation
from pptx.util import Pt

def rgb_to_hex(rgb_color):
    """Convert RGB color to hex"""
    try:
        return f"#{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
    except:
        return None

def get_master_background(pptx_path, slide_num):
    """Get the master background for a specific slide"""
    prs = Presentation(pptx_path)
    
    if slide_num >= len(prs.slides):
        print(f"Slide {slide_num} not found")
        return
    
    slide = prs.slides[slide_num]
    
    print(f"\n{'='*80}")
    print(f"MASTER BACKGROUND INSPECTION FOR SLIDE {slide_num + 1}")
    print(f"{'='*80}\n")
    
    # Get slide layout
    slide_layout = slide.slide_layout
    print(f"Slide Layout: {slide_layout.name}")
    
    # Check layout background
    print("\nLAYOUT BACKGROUND:")
    print("-" * 40)
    try:
        layout_bg = slide_layout.background
        layout_fill = layout_bg.fill
        print(f"Fill Type: {layout_fill.type}")
        
        if hasattr(layout_fill, 'fore_color'):
            try:
                rgb = layout_fill.fore_color.rgb
                hex_color = rgb_to_hex(rgb)
                print(f"Color (RGB): {rgb}")
                print(f"Color (Hex): {hex_color}")
            except:
                print("Could not resolve layout background color")
    except Exception as e:
        print(f"Layout background: {e}")
    
    # Get slide master
    slide_master = slide_layout.slide_master
    print(f"\nSlide Master: {slide_master.name if hasattr(slide_master, 'name') else 'Unknown'}")
    
    # Check master background
    print("\nMASTER BACKGROUND:")
    print("-" * 40)
    try:
        master_bg = slide_master.background
        master_fill = master_bg.fill
        print(f"Fill Type: {master_fill.type}")
        
        if hasattr(master_fill, 'fore_color'):
            try:
                if master_fill.fore_color.type == 1:  # RGB
                    rgb = master_fill.fore_color.rgb
                    hex_color = rgb_to_hex(rgb)
                    print(f"Color Type: RGB")
                    print(f"Color (RGB): {rgb}")
                    print(f"Color (Hex): {hex_color}")
                elif master_fill.fore_color.type == 2:  # Scheme
                    print(f"Color Type: Scheme")
                    print(f"Theme Color: {master_fill.fore_color.theme_color}")
                    try:
                        rgb = master_fill.fore_color.rgb
                        hex_color = rgb_to_hex(rgb)
                        print(f"Resolved (RGB): {rgb}")
                        print(f"Resolved (Hex): {hex_color}")
                    except:
                        print("Could not resolve scheme color")
            except Exception as e:
                print(f"Could not get master background color: {e}")
    except Exception as e:
        print(f"Master background: {e}")
    
    # Check theme colors
    print("\nTHEME COLORS:")
    print("-" * 40)
    try:
        theme = slide_master.theme
        if theme:
            color_scheme = theme.color_scheme
            colors = {
                'accent1': color_scheme.accent1,
                'accent2': color_scheme.accent2,
                'accent3': color_scheme.accent3,
                'dk1': color_scheme.dk1,
                'dk2': color_scheme.dk2,
                'lt1': color_scheme.lt1,
                'lt2': color_scheme.lt2,
            }
            
            for name, color in colors.items():
                try:
                    rgb = color.rgb
                    print(f"{name}: {rgb_to_hex(rgb)}")
                except:
                    print(f"{name}: Could not resolve")
    except Exception as e:
        print(f"Could not get theme colors: {e}")

if __name__ == "__main__":
    pptx_path = "branding/PwC_ppt_graphic_elements_level1.pptx"
    slide_num = 1  # Slide 2 (0-indexed)
    
    get_master_background(pptx_path, slide_num)




