#!/usr/bin/env python3
"""
Extract font properties from slide layout master
"""

from pptx import Presentation
from pathlib import Path
from lxml import etree

def get_layout_placeholder_properties(slide, shape):
    """Get properties from the slide layout for a placeholder"""
    
    if not shape.is_placeholder:
        return None
    
    try:
        # Get placeholder index
        ph_idx = shape.placeholder_format.idx
        
        # Get the slide layout
        layout = slide.slide_layout
        
        # Find the corresponding placeholder in the layout
        for layout_shape in layout.placeholders:
            if layout_shape.placeholder_format.idx == ph_idx:
                print(f"\nFound matching placeholder in layout, idx={ph_idx}")
                
                # Get the XML of this layout placeholder
                if hasattr(layout_shape, '_element'):
                    xml_str = etree.tostring(layout_shape._element, pretty_print=True, encoding='unicode')
                    print("Layout placeholder XML:")
                    print(xml_str[:1500])
                
                # Try to access text frame and get default properties
                if layout_shape.has_text_frame:
                    tf = layout_shape.text_frame
                    if hasattr(tf, '_element'):
                        # Look for lvl1pPr (level 1 paragraph properties) in lstStyle
                        lstStyle = tf._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
                        if lstStyle is not None:
                            print("\nFound lstStyle!")
                            lst_xml = etree.tostring(lstStyle, pretty_print=True, encoding='unicode')
                            print(lst_xml)
                
                break
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
    
    return None

def main():
    pptx_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_graphic_elements_level1.pptx")
    prs = Presentation(str(pptx_file))
    
    # Get slide 2
    slide = prs.slides[1]
    
    # Get shape 1 (the large "1")
    shape = slide.shapes[1]
    
    print(f"Shape text: '{shape.text}'")
    print(f"Is placeholder: {shape.is_placeholder}")
    
    get_layout_placeholder_properties(slide, shape)

if __name__ == "__main__":
    main()




