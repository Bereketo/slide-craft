#!/usr/bin/env python3
"""
Extract ALL images from all 75 covers without filtering.
This will help us see what images are actually in each cover.
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def main():
    covers_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_covers.pptx")
    covers_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/covers")
    images_dir = covers_dir / "images_all"
    
    images_dir.mkdir(exist_ok=True)
    
    print(f"Loading presentation: {covers_file}")
    prs = Presentation(str(covers_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides\n")
    
    all_images = {}
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"\nSlide {idx}:")
        slide_images = []
        
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Get dimensions
                    left_px = int(shape.left / 9525)
                    top_px = int(shape.top / 9525)
                    width_px = int(shape.width / 9525)
                    height_px = int(shape.height / 9525)
                    
                    # Get actual dimensions
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        actual_width, actual_height = img.size
                    except:
                        actual_width, actual_height = width_px, height_px
                    
                    # Determine file extension
                    ext_map = {
                        'image/jpeg': 'jpg',
                        'image/jpg': 'jpg',
                        'image/png': 'png',
                        'image/gif': 'gif'
                    }
                    ext = ext_map.get(image.content_type, 'jpg')
                    
                    # Save image
                    image_filename = f"cover_{idx:02d}_img{image_count}.{ext}"
                    image_path = images_dir / image_filename
                    
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    
                    info = {
                        "filename": image_filename,
                        "position": {"x": left_px, "y": top_px, "w": width_px, "h": height_px},
                        "actual_size": {"width": actual_width, "height": actual_height},
                        "content_type": image.content_type
                    }
                    
                    slide_images.append(info)
                    
                    print(f"  Image {image_count}: {image_filename}")
                    print(f"    Position: ({left_px}, {top_px}) Size: {width_px}x{height_px}")
                    print(f"    Actual: {actual_width}x{actual_height}")
                    
                except Exception as e:
                    print(f"    Error extracting image {image_count}: {e}")
        
        if slide_images:
            all_images[f"cover_{idx:02d}"] = slide_images
        else:
            print(f"  No images found")
    
    # Save index
    index_file = images_dir / "all_images_index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(all_images, f, indent=2, ensure_ascii=False)
    
    total_images = sum(len(imgs) for imgs in all_images.values())
    print(f"\n✓ Extraction complete!")
    print(f"  Total slides: {total_slides}")
    print(f"  Slides with images: {len(all_images)}")
    print(f"  Total images extracted: {total_images}")
    print(f"  Output directory: {images_dir}")

if __name__ == "__main__":
    main()








