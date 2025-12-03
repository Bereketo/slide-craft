#!/usr/bin/env python3
"""
Extract a single slide with correct background color
"""

import sys
import json
from pathlib import Path

# Import from extract_content_elements
sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from extras.extract_content_elements import extract_slide_to_json, extract_theme_colors_from_pptx

def extract_slide(pptx_path, slide_num, output_file):
    """Extract a specific slide"""
    # Extract theme colors first
    theme_colors = extract_theme_colors_from_pptx(pptx_path)
    if theme_colors:
        print(f"✓ Loaded theme colors: PwC Orange = {theme_colors.get('accent1', 'N/A')}")
    
    prs = Presentation(pptx_path)
    
    if slide_num >= len(prs.slides):
        print(f"Slide {slide_num} not found")
        return
    
    slide = prs.slides[slide_num]
    
    # Extract slide data with theme colors
    slide_data = extract_slide_to_json(slide, slide_num + 1, theme_colors)
    
    # Save to file
    with open(output_file, 'w') as f:
        json.dump(slide_data, f, indent=2)
    
    print(f"✓ Extracted slide {slide_num + 1} to {output_file}")
    print(f"  Background color: {slide_data['slide']['background']['color']}")
    print(f"  Components: {len(slide_data['slide']['components'])}")

if __name__ == "__main__":
    pptx_path = "branding/PwC_ppt_graphic_elements_level1.pptx"
    slide_num = 1  # Slide 2 (0-indexed)
    output_file = "branding/content_elements/title_slides/element_02_chart_corrected.json"
    
    extract_slide(pptx_path, slide_num, output_file)

