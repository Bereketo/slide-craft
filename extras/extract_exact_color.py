#!/usr/bin/env python3
"""
Extract exact color from the '1' text in the original slide
"""

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE

def rgb_to_hex(rgb_color):
    """Convert RGB color to hex"""
    return f"#{str(rgb_color)}"

def extract_text_color(pptx_path, slide_num):
    """Extract color from specific text"""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num]
    
    print(f"\n{'='*80}")
    print(f"EXTRACTING COLOR FROM SLIDE {slide_num + 1}")
    print(f"{'='*80}\n")
    
    for idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        
        text = shape.text_frame.text.strip()
        if text == "1":
            print(f"Found '1' in shape {idx + 1}")
            print(f"Position: x={shape.left // 9525}px, y={shape.top // 9525}px\n")
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    print(f"Text: '{run.text}'")
                    print(f"Font: {run.font.name}")
                    print(f"Size: {run.font.size.pt if run.font.size else 'Default'}pt")
                    print(f"Bold: {run.font.bold}")
                    
                    # Get color
                    try:
                        color_obj = run.font.color
                        print(f"\nColor Type: {color_obj.type}")
                        
                        if color_obj.type == MSO_COLOR_TYPE.RGB:
                            rgb = color_obj.rgb
                            hex_color = rgb_to_hex(rgb)
                            print(f"Color (RGB): {rgb}")
                            print(f"Color (Hex): {hex_color}")
                            print(f"Color (repr): {repr(rgb)}")
                            
                        elif color_obj.type == MSO_COLOR_TYPE.SCHEME:
                            print(f"Color is SCHEME: {color_obj.theme_color}")
                            try:
                                rgb = color_obj.rgb
                                hex_color = rgb_to_hex(rgb)
                                print(f"Resolved RGB: {rgb}")
                                print(f"Resolved Hex: {hex_color}")
                                print(f"Resolved (repr): {repr(rgb)}")
                            except Exception as e:
                                print(f"Could not resolve: {e}")
                        
                        else:
                            print(f"Unknown color type: {color_obj.type}")
                    
                    except Exception as e:
                        print(f"Error getting color: {e}")
                    
                    break  # Just first run
                break  # Just first paragraph

if __name__ == "__main__":
    pptx_path = "branding/PwC_ppt_graphic_elements_level1.pptx"
    slide_num = 1  # Slide 2 (0-indexed)
    
    extract_text_color(pptx_path, slide_num)




