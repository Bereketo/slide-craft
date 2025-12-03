#!/usr/bin/env python3
"""
Extract PwC content elements from the graphic elements template PPTX and convert to JSON format.
Each element will be saved as a separate JSON file with metadata about its use case.
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
import xml.etree.ElementTree as ET

def extract_theme_colors_from_pptx(pptx_path):
    """Extract theme colors dynamically from the PPTX file"""
    theme_colors = {}
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_file:
            # Get all theme files
            theme_files = [name for name in zip_file.namelist() if 'theme/theme1.xml' in name]
            
            if theme_files:
                with zip_file.open(theme_files[0]) as f:
                    content = f.read()
                    root = ET.fromstring(content)
                    
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    clrScheme = root.find('.//a:clrScheme', ns)
                    
                    if clrScheme:
                        for color_elem in clrScheme:
                            color_name = color_elem.tag.split('}')[-1]
                            
                            # Look for srgbClr
                            srgb = color_elem.find('.//a:srgbClr', ns)
                            if srgb is not None:
                                val = srgb.get('val')
                                theme_colors[color_name] = f"#{val}"
                            
                            # Look for sysClr
                            sys = color_elem.find('.//a:sysClr', ns)
                            if sys is not None:
                                val = sys.get('lastClr')
                                if val:
                                    theme_colors[color_name] = f"#{val}"
    except Exception as e:
        print(f"Warning: Could not extract theme colors: {e}")
    
    return theme_colors

def emu_to_px(emu_value):
    """Convert EMU to pixels (96 DPI)"""
    return int(emu_value / 9525)

def rgb_to_hex(rgb_color):
    """Convert RGB color to hex"""
    try:
        # RGBColor can be converted to string directly (returns hex without #)
        # Or indexed like rgb[0], rgb[1], rgb[2]
        return f"#{str(rgb_color)}"
    except:
        try:
            # Fallback: try indexing
            return f"#{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
        except:
            return "#000000"

def get_fill_color(shape):
    """Extract fill color from shape"""
    try:
        if shape.fill.type == 1:  # SOLID
            return rgb_to_hex(shape.fill.fore_color.rgb)
        elif shape.fill.type == 0:  # No fill
            return None
    except:
        pass
    return None

def get_inherited_properties(slide, shape, paragraph, run, theme_colors=None):
    """Get inherited font properties from slide layout when properties are None"""
    props = {
        "font_size": None,
        "font_family": None,
        "color": None,
        "bold": None
    }
    
    # Use dynamically loaded theme colors from theme1.xml, or fallback to defaults
    SCHEME_COLORS = theme_colors if theme_colors else {
        "accent1": "#FD5108",  # PwC Orange (from theme1.xml)
        "accent2": "#FE7C39",  
        "accent3": "#FFAA72",
        "dk1": "#000000",      # Black
        "dk2": "#000000",
        "lt1": "#FFFFFF",      # White
        "lt2": "#EBEBEB",      # Light grey
        "tx1": "#000000",      # Black text
        "tx2": "#595959",      # Gray text
        "bg1": "#FFFFFF",      # White background
        "bg2": "#F2F2F2",      # Light gray background
    }
    
    # If it's a placeholder, get properties from slide layout
    if shape.is_placeholder:
        try:
            ph_idx = shape.placeholder_format.idx
            layout = slide.slide_layout
            
            # Find matching placeholder in layout
            for layout_shape in layout.placeholders:
                if layout_shape.placeholder_format.idx == ph_idx:
                    if layout_shape.has_text_frame:
                        tf = layout_shape.text_frame
                        if hasattr(tf, '_element'):
                            # Get lstStyle
                            lstStyle = tf._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
                            if lstStyle is not None:
                                # Get lvl1pPr (level 1 paragraph properties)
                                lvl1pPr = lstStyle.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
                                if lvl1pPr is not None:
                                    # Get defRPr (default run properties)
                                    defRPr = lvl1pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
                                    if defRPr is not None:
                                        # Font size (sz is in hundredths of a point, e.g., 1050 = 10.5pt)
                                        sz = defRPr.get('sz')
                                        if sz:
                                            props["font_size"] = int(sz) / 100  # Keep decimal precision
                                        
                                        # Bold
                                        b = defRPr.get('b')
                                        if b is not None:
                                            props["bold"] = (b == '1' or b.lower() == 'true')
                                        
                                        # Font family
                                        latin = defRPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                                        if latin is not None:
                                            typeface = latin.get('typeface')
                                            if typeface and typeface != '+mj-lt' and typeface != '+mn-lt':
                                                props["font_family"] = typeface
                                        
                                        # Color
                                        solidFill = defRPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                                        if solidFill is not None:
                                            # Check for scheme color first
                                            schemeClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
                                            if schemeClr is not None:
                                                scheme_val = schemeClr.get('val')
                                                if scheme_val and scheme_val in SCHEME_COLORS:
                                                    props["color"] = SCHEME_COLORS[scheme_val]
                                            else:
                                                # Check for direct RGB color
                                                srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                                                if srgbClr is not None:
                                                    color_val = srgbClr.get('val')
                                                    if color_val:
                                                        props["color"] = f"#{color_val.upper()}"
                    break
        except Exception as e:
            pass
    
    return props

def extract_text_from_shape(slide, shape, theme_colors=None):
    """Extract text and styling from a shape"""
    if not shape.has_text_frame:
        return None
    
    text_frame = shape.text_frame
    if not text_frame.text.strip():
        return None
    
    components = []
    
    for paragraph in text_frame.paragraphs:
        if not paragraph.text.strip():
            continue
            
        runs = []
        for run in paragraph.runs:
            if not run.text:
                continue
            
            # Get inherited properties for placeholders
            inherited = get_inherited_properties(slide, shape, paragraph, run, theme_colors)
                
            # Extract font size - use run property or inherited (keep as float for precision)
            font_size = 18  # default
            try:
                if run.font.size:
                    font_size = run.font.size.pt  # Keep as float, don't convert to int
                elif inherited["font_size"]:
                    font_size = inherited["font_size"]
            except:
                if inherited["font_size"]:
                    font_size = inherited["font_size"]
            
            # Get bold
            bold = False
            if run.font.bold is not None:
                bold = run.font.bold
            elif inherited["bold"] is not None:
                bold = inherited["bold"]
            
            run_info = {
                "text": run.text,
                "bold": bold,
                "font_size": font_size,
            }
            
            # Try to get font family - use run property or inherited
            font_family = None
            try:
                if run.font.name:
                    font_family = run.font.name
                elif inherited["font_family"]:
                    font_family = inherited["font_family"]
            except:
                if inherited["font_family"]:
                    font_family = inherited["font_family"]
            
            if font_family:
                run_info["font_family"] = font_family
            
            # Try to get color - use run property or inherited
            color = None
            try:
                color_obj = run.font.color
                if color_obj.type == 1:  # RGB (MSO_COLOR_TYPE.RGB)
                    color = rgb_to_hex(color_obj.rgb)
                elif color_obj.type == 2:  # SCHEME (MSO_COLOR_TYPE.SCHEME)
                    try:
                        color = rgb_to_hex(color_obj.rgb)
                    except:
                        pass
                elif hasattr(color_obj, 'rgb'):
                    color = rgb_to_hex(color_obj.rgb)
            except:
                pass
            
            if not color and inherited["color"]:
                color = inherited["color"]
            
            if color:
                run_info["color"] = color
                
            runs.append(run_info)
        
        if runs:
            alignment = "left"
            try:
                if paragraph.alignment:
                    alignment = str(paragraph.alignment).split('.')[-1].lower()
                    # Clean up alignment values
                    if 'center' in alignment:
                        alignment = 'center'
                    elif 'right' in alignment:
                        alignment = 'right'
                    elif 'justify' in alignment:
                        alignment = 'justify'
                    else:
                        alignment = 'left'
            except:
                pass
            
            components.append({
                "runs": runs,
                "alignment": alignment
            })
    
    return components if components else None

def categorize_element(slide_num, slide, text_content):
    """Categorize the element based on its content and layout"""
    
    text_lower = text_content.lower()
    shape_types = [s.shape_type for s in slide.shapes]
    shape_count = len(slide.shapes)
    
    # Analyze content to determine category
    if "chart" in text_lower or "graph" in text_lower:
        return "chart", "Chart or graph layout"
    elif "table" in text_lower or "column" in text_lower or "row" in text_lower:
        return "table", "Table or data layout"
    elif "bullet" in text_lower or "list" in text_lower:
        return "bullet_list", "Bullet point or list layout"
    elif "title" in text_lower and "content" in text_lower:
        return "title_content", "Title and content layout"
    elif "image" in text_lower or "photo" in text_lower:
        return "image_layout", "Image-focused layout"
    elif "two" in text_lower or "split" in text_lower or "column" in text_lower:
        return "two_column", "Two-column layout"
    elif "section" in text_lower or "divider" in text_lower:
        return "section_divider", "Section divider"
    elif "blank" in text_lower or shape_count < 3:
        return "blank", "Blank or minimal layout"
    elif "quote" in text_lower or "testimonial" in text_lower:
        return "quote", "Quote or testimonial layout"
    elif "timeline" in text_lower or "process" in text_lower:
        return "timeline", "Timeline or process layout"
    elif "comparison" in text_lower or "vs" in text_lower:
        return "comparison", "Comparison layout"
    else:
        # Default categorization based on shape count
        if shape_count > 10:
            return "complex", "Complex multi-element layout"
        elif shape_count > 5:
            return "content", "Standard content layout"
        else:
            return "simple", "Simple layout"

def get_slide_background_color(slide):
    """Get the actual background color, checking layout if needed"""
    try:
        # Check slide background first
        background = slide.background
        fill = background.fill
        
        if fill.type == 1:  # Solid fill
            try:
                rgb = fill.fore_color.rgb
                return rgb_to_hex(rgb)
            except:
                pass
        
        # If slide background is inherited (type 5), check layout
        if fill.type == 5:  # BACKGROUND type - inherited
            try:
                layout = slide.slide_layout
                layout_fill = layout.background.fill
                
                if layout_fill.type == 1:  # Solid
                    rgb = layout_fill.fore_color.rgb
                    return rgb_to_hex(rgb)
            except:
                pass
    except:
        pass
    
    return "#FFFFFF"  # Default to white

def extract_slide_to_json(slide, slide_num, theme_colors=None):
    """Extract a single slide and convert to our JSON schema format"""
    
    components = []
    all_text = ""
    
    # Extract all shapes
    for shape in slide.shapes:
        # Skip placeholders that are empty
        if shape.is_placeholder and not shape.has_text_frame:
            continue
        
        # Get position and size
        left_px = emu_to_px(shape.left)
        top_px = emu_to_px(shape.top)
        width_px = emu_to_px(shape.width)
        height_px = emu_to_px(shape.height)
        
        # Extract text components
        if shape.has_text_frame:
            text_components = extract_text_from_shape(slide, shape, theme_colors)
            if text_components:
                all_text += " " + shape.text
                
                for text_comp in text_components:
                    # Determine text type based on font size
                    avg_size = sum(r.get("font_size", 18) for r in text_comp["runs"]) / len(text_comp["runs"])
                    
                    if avg_size >= 36:
                        text_type = "title"
                    elif avg_size >= 24:
                        text_type = "h2"
                    else:
                        text_type = "body"
                    
                    component = {
                        "type": "richtext",
                        "runs": text_comp["runs"],
                        "box": {
                            "x": left_px,
                            "y": top_px,
                            "w": width_px,
                            "h": height_px,
                            "unit": "px"
                        },
                        "style": {
                            "align": text_comp["alignment"]
                        }
                    }
                    
                    components.append(component)
        
        # Extract image components
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            component = {
                "type": "image",
                "src": f"images/element_{slide_num:02d}_placeholder.png",
                "alt": f"Element image {slide_num}",
                "box": {
                    "x": left_px,
                    "y": top_px,
                    "w": width_px,
                    "h": height_px,
                    "unit": "px"
                },
                "object_fit": "cover"
            }
            components.append(component)
        
        # Extract shape components
        elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
            fill_color = get_fill_color(shape)
            if fill_color:
                component = {
                    "type": "shape",
                    "shape_type": "rectangle",
                    "box": {
                        "x": left_px,
                        "y": top_px,
                        "w": width_px,
                        "h": height_px,
                        "unit": "px"
                    },
                    "style": {
                        "fill": fill_color
                    }
                }
                components.append(component)
    
    # Categorize the element
    category, description = categorize_element(slide_num, slide, all_text)
    
    # Get the actual background color
    bg_color = get_slide_background_color(slide)
    
    return {
        "metadata": {
            "element_number": slide_num,
            "category": category,
            "description": description,
            "extracted_text": all_text.strip()[:200]  # First 200 chars for reference
        },
        "slide": {
            "title": f"PwC Content Element {slide_num}",
            "background": {"type": "solid", "color": bg_color},
            "components": components
        }
    }

def main():
    # Path to the PwC graphic elements file
    elements_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_graphic_elements_level1.pptx")
    
    if not elements_file.exists():
        print(f"Error: File not found: {elements_file}")
        return
    
    # Extract theme colors from the PPTX file
    print(f"Extracting theme colors from: {elements_file}")
    theme_colors = extract_theme_colors_from_pptx(str(elements_file))
    if theme_colors:
        print(f"✓ Extracted {len(theme_colors)} theme colors")
        print(f"  PwC Orange (accent1): {theme_colors.get('accent1', 'Not found')}")
    else:
        print("⚠ Warning: Could not extract theme colors, using defaults")
    
    # Create output directory
    output_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/content_elements")
    output_dir.mkdir(exist_ok=True)
    
    # Create subdirectories
    elements_dir = output_dir / "elements"
    elements_dir.mkdir(exist_ok=True)
    
    print(f"\nLoading presentation from: {elements_file}")
    prs = Presentation(str(elements_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides in presentation")
    
    # Extract each slide
    all_elements = []
    category_counts = {}
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {idx}/{total_slides}...")
        
        try:
            element_json = extract_slide_to_json(slide, idx, theme_colors)
            
            # Save individual element
            output_file = elements_dir / f"element_{idx:02d}_{element_json['metadata']['category']}.json"
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(element_json, f, indent=2, ensure_ascii=False)
            
            all_elements.append(element_json)
            
            # Track categories
            category = element_json['metadata']['category']
            category_counts[category] = category_counts.get(category, 0) + 1
            
            print(f"  ✓ Saved: {output_file.name}")
            
        except Exception as e:
            print(f"  ✗ Error processing slide {idx}: {e}")
            import traceback
            traceback.print_exc()
    
    # Create index file
    index = {
        "total_elements": len(all_elements),
        "categories": category_counts,
        "elements": [
            {
                "id": element["metadata"]["element_number"],
                "file": f"element_{element['metadata']['element_number']:02d}_{element['metadata']['category']}.json",
                "category": element["metadata"]["category"],
                "description": element["metadata"]["description"]
            }
            for element in all_elements
        ]
    }
    
    index_file = output_dir / "index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(index, f, indent=2, ensure_ascii=False)
    
    print(f"\n✓ Extraction complete!")
    print(f"  Total elements extracted: {len(all_elements)}")
    print(f"  Output directory: {output_dir}")
    print(f"  Index file: {index_file}")
    print(f"\nCategory breakdown:")
    for category, count in sorted(category_counts.items()):
        print(f"  - {category}: {count} elements")

if __name__ == "__main__":
    main()

