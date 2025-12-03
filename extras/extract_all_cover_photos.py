#!/usr/bin/env python3
"""
Extract all 75 unique cover photos from PwC covers.
- Logo goes to logo/ directory
- Each cover's unique photo goes to images/ directory
"""

import json
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def is_logo_image(width, height, position_x, position_y):
    """
    Determine if an image is likely a logo vs a cover photo.
    Logo is typically:
    - Smaller (< 200px width or height)
    - Positioned at top-left
    - Aspect ratio close to 2:1 for PwC logo
    """
    # Small image at top-left = logo
    if width < 200 and height < 100 and position_x < 200 and position_y < 100:
        return True
    return False

def extract_cover_photo(slide, slide_num, output_dir):
    """Extract the main cover photo (not the logo) from a slide"""
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                
                # Get dimensions
                left_px = int(shape.left / 9525)
                top_px = int(shape.top / 9525)
                width_px = int(shape.width / 9525)
                height_px = int(shape.height / 9525)
                
                # Skip if it's the logo
                if is_logo_image(width_px, height_px, left_px, top_px):
                    continue
                
                # This is the cover photo - save it
                ext_map = {
                    'image/jpeg': 'jpg',
                    'image/jpg': 'jpg',
                    'image/png': 'png',
                    'image/gif': 'gif'
                }
                ext = ext_map.get(image.content_type, 'jpg')
                
                image_filename = f"cover_{slide_num:02d}_photo.{ext}"
                image_path = output_dir / image_filename
                
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                # Get actual dimensions
                try:
                    img = Image.open(io.BytesIO(image_bytes))
                    actual_width, actual_height = img.size
                except:
                    actual_width, actual_height = width_px, height_px
                
                return {
                    "filename": image_filename,
                    "position": {"x": left_px, "y": top_px, "w": width_px, "h": height_px},
                    "actual_dimensions": {"width": actual_width, "height": actual_height},
                    "content_type": image.content_type
                }
                
            except Exception as e:
                print(f"    Error: {e}")
                continue
    
    return None

def main():
    covers_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_covers.pptx")
    covers_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/covers")
    images_dir = covers_dir / "images"
    
    # Clean and recreate images directory
    if images_dir.exists():
        shutil.rmtree(images_dir)
    images_dir.mkdir(exist_ok=True)
    
    print(f"Loading presentation: {covers_file}")
    prs = Presentation(str(covers_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides\n")
    
    extracted_count = 0
    cover_photos = {}
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {idx}/{total_slides}...", end=" ")
        
        try:
            photo_info = extract_cover_photo(slide, idx, images_dir)
            
            if photo_info:
                cover_photos[idx] = photo_info
                extracted_count += 1
                print(f"✓ {photo_info['filename']} ({photo_info['actual_dimensions']['width']}x{photo_info['actual_dimensions']['height']})")
                
                # Update JSON file
                json_file = covers_dir / f"cover_{idx:02d}_title_subtitle.json"
                if json_file.exists():
                    with open(json_file, 'r', encoding='utf-8') as f:
                        cover_data = json.load(f)
                    
                    # Update components to reference the cover photo
                    for component in cover_data['slide']['components']:
                        if component['type'] == 'image':
                            # Check if it's likely the cover photo (not logo)
                            if 'box' in component:
                                box = component['box']
                                if not is_logo_image(box['w'], box['h'], box['x'], box['y']):
                                    component['src'] = f"images/{photo_info['filename']}"
                                    component['image_metadata'] = {
                                        "actual_width": photo_info['actual_dimensions']['width'],
                                        "actual_height": photo_info['actual_dimensions']['height'],
                                        "content_type": photo_info['content_type']
                                    }
                                else:
                                    component['src'] = "logo/pwc_logo.png"
                    
                    with open(json_file, 'w', encoding='utf-8') as f:
                        json.dump(cover_data, f, indent=2, ensure_ascii=False)
            else:
                print("- No cover photo found")
                
        except Exception as e:
            print(f"✗ Error: {e}")
    
    # Create cover photos index
    index = {
        "total_covers": total_slides,
        "covers_with_photos": extracted_count,
        "covers": {
            f"cover_{idx:02d}": info
            for idx, info in cover_photos.items()
        }
    }
    
    index_file = images_dir / "covers_index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(index, f, indent=2, ensure_ascii=False)
    
    print(f"\n✓ Extraction complete!")
    print(f"  Total covers: {total_slides}")
    print(f"  Cover photos extracted: {extracted_count}")
    print(f"  Output directory: {images_dir}")
    print(f"  Index file: {index_file}")

if __name__ == "__main__":
    main()








