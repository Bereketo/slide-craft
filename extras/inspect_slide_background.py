#!/usr/bin/env python3
"""
Inspect slide background colors and exact positioning from PwC template
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR

def rgb_to_hex(rgb_color):
    """Convert RGB color to hex"""
    try:
        return f"#{rgb_color[0]:02X}{rgb_color[1]:02X}{rgb_color[2]:02X}"
    except:
        return None

def inspect_slide(pptx_path, slide_num):
    """Inspect specific slide"""
    prs = Presentation(pptx_path)
    
    if slide_num >= len(prs.slides):
        print(f"Slide {slide_num} not found")
        return
    
    slide = prs.slides[slide_num]
    
    print(f"\n{'='*80}")
    print(f"SLIDE {slide_num + 1} INSPECTION")
    print(f"{'='*80}\n")
    
    # Check background
    print("BACKGROUND:")
    print("-" * 40)
    
    try:
        background = slide.background
        fill = background.fill
        
        if fill.type == 1:  # Solid fill
            print("Type: Solid Fill")
            try:
                if fill.fore_color.type == 1:  # RGB
                    rgb = fill.fore_color.rgb
                    hex_color = rgb_to_hex(rgb)
                    print(f"Color (RGB): {rgb}")
                    print(f"Color (Hex): {hex_color}")
                elif fill.fore_color.type == 2:  # Scheme
                    print(f"Color Type: Scheme")
                    print(f"Scheme Color: {fill.fore_color.theme_color}")
                    try:
                        rgb = fill.fore_color.rgb
                        hex_color = rgb_to_hex(rgb)
                        print(f"Resolved RGB: {rgb}")
                        print(f"Resolved Hex: {hex_color}")
                    except:
                        print("Could not resolve scheme color to RGB")
            except Exception as e:
                print(f"Could not get fill color: {e}")
        else:
            print(f"Fill type: {fill.type}")
    except Exception as e:
        print(f"Could not inspect background: {e}")
        print("Background might be inherited from master")
    
    # Check shapes
    print(f"\n\nSHAPES ({len(slide.shapes)} total):")
    print("-" * 40)
    
    for idx, shape in enumerate(slide.shapes):
        print(f"\n[Shape {idx + 1}]")
        print(f"Type: {shape.shape_type}")
        print(f"Position: x={shape.left}, y={shape.top}")
        print(f"Size: w={shape.width}, h={shape.height}")
        print(f"Position (px): x={shape.left // 9525}, y={shape.top // 9525}")
        print(f"Size (px): w={shape.width // 9525}, h={shape.height // 9525}")
        
        if shape.has_text_frame:
            text = shape.text_frame.text
            print(f"Text: {text[:50]}...")
            
            # Check shape fill
            if shape.fill.type == 1:  # Solid
                try:
                    rgb = shape.fill.fore_color.rgb
                    print(f"Shape Fill: {rgb_to_hex(rgb)}")
                except:
                    print("Shape Fill: Could not determine")
            
            # Check text formatting
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    print(f"  Font: {run.font.name if run.font.name else 'Default'}")
                    print(f"  Size: {run.font.size.pt if run.font.size else 'Default'}pt")
                    try:
                        if run.font.color.type == 1:
                            rgb = run.font.color.rgb
                            print(f"  Color: {rgb_to_hex(rgb)}")
                        elif run.font.color.type == 2:
                            print(f"  Color: Scheme ({run.font.color.theme_color})")
                            try:
                                rgb = run.font.color.rgb
                                print(f"  Color (resolved): {rgb_to_hex(rgb)}")
                            except:
                                pass
                    except:
                        pass
                    break  # Just first run
                break  # Just first paragraph

if __name__ == "__main__":
    pptx_path = "branding/PwC_ppt_graphic_elements_level1.pptx"
    slide_num = 1  # 0-indexed, so this is slide 2
    
    inspect_slide(pptx_path, slide_num)




