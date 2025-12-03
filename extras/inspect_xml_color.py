#!/usr/bin/env python3
"""
Inspect XML directly to get the exact color value
"""

from pptx import Presentation
import xml.etree.ElementTree as ET

def inspect_xml_color(pptx_path, slide_num):
    """Inspect XML for color information"""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num]
    
    # Find the shape with "1"
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "1":
            print("Found shape with '1'")
            print(f"Shape type: {shape.shape_type}")
            
            # Get the XML
            xml = shape._element.xml
            
            # Parse and pretty print relevant parts
            root = ET.fromstring(xml)
            
            # Look for color elements
            print("\nSearching for color elements in XML...")
            
            # Find all solidFill elements
            for elem in root.iter():
                if 'solidFill' in elem.tag:
                    print(f"\nFound solidFill: {elem.tag}")
                    for child in elem:
                        print(f"  Child: {child.tag}, attribs: {child.attrib}")
                        if 'schemeClr' in child.tag:
                            print(f"    Scheme Color: {child.get('val')}")
                        elif 'srgbClr' in child.tag:
                            print(f"    RGB Color: {child.get('val')}")
                
                if 'schemeClr' in elem.tag:
                    print(f"\nFound schemeClr directly: val={elem.get('val')}")
                
                if 'srgbClr' in elem.tag:
                    print(f"\nFound srgbClr directly: val={elem.get('val')}")
            
            break

if __name__ == "__main__":
    pptx_path = "branding/PwC_ppt_graphic_elements_level1.pptx"
    slide_num = 1
    
    inspect_xml_color(pptx_path, slide_num)




