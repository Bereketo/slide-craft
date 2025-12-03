#!/usr/bin/env python3
"""
Deep XML inspection to find where font properties are stored
"""

from pptx import Presentation
from pathlib import Path
from lxml import etree

def inspect_xml_structure():
    pptx_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_graphic_elements_level1.pptx")
    prs = Presentation(str(pptx_file))
    
    # Get slide 2 (the "1" with orange color and size 350)
    slide = prs.slides[1]
    
    # Get shape 1 (the large "1")
    shape = slide.shapes[1]
    
    print("=== SHAPE 1 (The large '1') ===\n")
    
    # Get the text frame
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]
    
    print("Text:", run.text)
    print("\n=== RUN XML ===")
    if hasattr(run, '_element'):
        print(etree.tostring(run._element, pretty_print=True, encoding='unicode'))
    
    print("\n=== PARAGRAPH XML ===")
    if hasattr(paragraph, '_element'):
        print(etree.tostring(paragraph._element, pretty_print=True, encoding='unicode'))
    
    print("\n=== TEXT FRAME XML ===")
    if hasattr(text_frame, '_element'):
        print(etree.tostring(text_frame._element, pretty_print=True, encoding='unicode'))
    
    print("\n=== SHAPE XML (first 2000 chars) ===")
    if hasattr(shape, '_element'):
        xml_str = etree.tostring(shape._element, pretty_print=True, encoding='unicode')
        print(xml_str[:2000])
    
    # Try to access slide layout
    print("\n=== CHECKING SLIDE LAYOUT ===")
    if hasattr(slide, 'slide_layout'):
        layout = slide.slide_layout
        print(f"Layout name: {layout.name if hasattr(layout, 'name') else 'Unknown'}")
        
        # Try to find corresponding placeholder in layout
        if shape.is_placeholder:
            print(f"Shape is placeholder, idx: {shape.placeholder_format.idx if hasattr(shape, 'placeholder_format') else 'Unknown'}")

if __name__ == "__main__":
    inspect_xml_structure()




