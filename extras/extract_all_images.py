#!/usr/bin/env python3
"""
Extract ALL images from a PowerPoint file and save them to an images directory.
This script extracts every image from every slide, including logos, photos, and graphics.
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def extract_all_images_from_slide(slide, slide_num, output_dir):
    """Extract all images from a slide and save them with metadata."""
    
    images_info = []
    
    # Get all PICTURE shapes from the slide
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    
    if not picture_shapes:
        return images_info
    
    for idx, shape in enumerate(picture_shapes, start=1):
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Map content types to file extensions
            ext_map = {
                'image/jpeg': 'jpg',
                'image/jpg': 'jpg',
                'image/png': 'png',
                'image/gif': 'gif',
                'image/bmp': 'bmp',
                'image/tiff': 'tiff',
                'image/svg+xml': 'svg'
            }
            ext = ext_map.get(image.content_type, 'png')
            
            # Create filename
            image_filename = f"slide_{slide_num:02d}_image_{idx}.{ext}"
            image_path = output_dir / image_filename
            
            # Save the image
            with open(image_path, 'wb') as f:
                f.write(image_bytes)
            
            # Get position and size in pixels (EMU to pixels conversion: divide by 9525)
            left_px = int(shape.left / 9525)
            top_px = int(shape.top / 9525)
            width_px = int(shape.width / 9525)
            height_px = int(shape.height / 9525)
            
            # Get actual image dimensions
            try:
                img = Image.open(io.BytesIO(image_bytes))
                actual_width, actual_height = img.size
            except:
                actual_width, actual_height = width_px, height_px
            
            # Store metadata
            image_info = {
                "filename": image_filename,
                "image_number": idx,
                "position": {
                    "x": left_px,
                    "y": top_px,
                    "width": width_px,
                    "height": height_px,
                    "unit": "px"
                },
                "actual_dimensions": {
                    "width": actual_width,
                    "height": actual_height
                },
                "content_type": image.content_type,
                "file_size_bytes": len(image_bytes)
            }
            
            images_info.append(image_info)
            print(f"    ✓ Extracted: {image_filename} ({actual_width}x{actual_height}) - {len(image_bytes)} bytes")
            
        except Exception as e:
            print(f"    ✗ Error extracting image {idx} from slide {slide_num}: {e}")
    
    return images_info

def main():
    # Define paths
    pptx_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_covers.pptx")
    output_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/images")
    
    # Create output directory
    output_dir.mkdir(exist_ok=True)
    
    # Check if file exists
    if not pptx_file.exists():
        print(f"Error: File not found: {pptx_file}")
        return
    
    print(f"Loading presentation from: {pptx_file}")
    print(f"Output directory: {output_dir}\n")
    
    # Load presentation
    prs = Presentation(str(pptx_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides in presentation\n")
    
    # Extract images from all slides
    all_extractions = {}
    total_images = 0
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {idx}/{total_slides}...")
        
        try:
            images_info = extract_all_images_from_slide(slide, idx, output_dir)
            
            if images_info:
                all_extractions[f"slide_{idx:02d}"] = {
                    "slide_number": idx,
                    "image_count": len(images_info),
                    "images": images_info
                }
                total_images += len(images_info)
            else:
                print(f"    - No images found")
                
        except Exception as e:
            print(f"  ✗ Error processing slide {idx}: {e}")
    
    # Create index file with all metadata
    image_index = {
        "source_file": str(pptx_file.name),
        "extraction_timestamp": "2025-11-04",
        "total_slides": total_slides,
        "slides_with_images": len(all_extractions),
        "total_images_extracted": total_images,
        "extractions": all_extractions
    }
    
    index_file = output_dir / "image_index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(image_index, f, indent=2, ensure_ascii=False)
    
    # Print summary
    print(f"\n{'='*60}")
    print(f"✓ Image extraction complete!")
    print(f"{'='*60}")
    print(f"  Source file: {pptx_file.name}")
    print(f"  Total slides: {total_slides}")
    print(f"  Slides with images: {len(all_extractions)}")
    print(f"  Total images extracted: {total_images}")
    print(f"  Output directory: {output_dir}")
    print(f"  Index file: {index_file}")
    print(f"{'='*60}\n")

if __name__ == "__main__":
    main()



