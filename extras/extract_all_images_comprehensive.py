#!/usr/bin/env python3
"""
Extract ALL images from a PowerPoint file including:
- Direct picture shapes
- Images in grouped shapes
- Images as shape fills
- Background images
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def extract_image_from_shape(shape, slide_num, image_counter, output_dir, source_type="direct"):
    """Extract an image from a shape and save it with metadata."""
    try:
        image = shape.image
        image_bytes = image.blob
        
        # Map content types to file extensions
        ext_map = {
            'image/jpeg': 'jpg',
            'image/jpg': 'jpg',
            'image/png': 'png',
            'image/gif': 'gif',
            'image/bmp': 'bmp',
            'image/tiff': 'tiff',
            'image/svg+xml': 'svg'
        }
        ext = ext_map.get(image.content_type, 'png')
        
        # Create filename
        image_filename = f"slide_{slide_num:02d}_image_{image_counter}.{ext}"
        image_path = output_dir / image_filename
        
        # Save the image
        with open(image_path, 'wb') as f:
            f.write(image_bytes)
        
        # Get position and size in pixels (EMU to pixels conversion: divide by 9525)
        left_px = int(shape.left / 9525) if hasattr(shape, 'left') else 0
        top_px = int(shape.top / 9525) if hasattr(shape, 'top') else 0
        width_px = int(shape.width / 9525) if hasattr(shape, 'width') else 0
        height_px = int(shape.height / 9525) if hasattr(shape, 'height') else 0
        
        # Get actual image dimensions
        try:
            img = Image.open(io.BytesIO(image_bytes))
            actual_width, actual_height = img.size
        except:
            actual_width, actual_height = width_px, height_px
        
        # Store metadata
        image_info = {
            "filename": image_filename,
            "image_number": image_counter,
            "source_type": source_type,
            "position": {
                "x": left_px,
                "y": top_px,
                "width": width_px,
                "height": height_px,
                "unit": "px"
            },
            "actual_dimensions": {
                "width": actual_width,
                "height": actual_height
            },
            "content_type": image.content_type,
            "file_size_bytes": len(image_bytes)
        }
        
        print(f"    ✓ Extracted [{source_type}]: {image_filename} ({actual_width}x{actual_height}) - {len(image_bytes)} bytes")
        return image_info
        
    except Exception as e:
        print(f"    ✗ Error extracting image: {e}")
        return None

def extract_all_images_from_slide(slide, slide_num, output_dir):
    """Extract all images from a slide from various sources."""
    
    images_info = []
    image_counter = 1
    
    # Check slide background for image fill
    if slide.background.fill.type == 6:  # MSO_FILL_TYPE.PICTURE
        try:
            fill = slide.background.fill
            if hasattr(fill, 'image'):
                class BackgroundShape:
                    def __init__(self, fill, slide):
                        self.image = fill.image
                        self.left = 0
                        self.top = 0
                        self.width = slide.shapes[0].width if slide.shapes else 1280 * 9525
                        self.height = slide.shapes[0].height if slide.shapes else 720 * 9525
                
                bg_shape = BackgroundShape(fill, slide)
                image_info = extract_image_from_shape(bg_shape, slide_num, image_counter, output_dir, "background")
                if image_info:
                    images_info.append(image_info)
                    image_counter += 1
        except Exception as e:
            print(f"    - Background check error: {e}")
    
    # Process all shapes
    for shape in slide.shapes:
        # Direct PICTURE shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_info = extract_image_from_shape(shape, slide_num, image_counter, output_dir, "picture")
            if image_info:
                images_info.append(image_info)
                image_counter += 1
        
        # GROUP shapes - recursively check for pictures
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_info = extract_image_from_shape(sub_shape, slide_num, image_counter, output_dir, "group")
                    if image_info:
                        images_info.append(image_info)
                        image_counter += 1
        
        # Shapes with picture fill
        elif hasattr(shape, 'fill') and shape.fill.type == 6:  # MSO_FILL_TYPE.PICTURE
            try:
                if hasattr(shape.fill, 'image'):
                    image_info = extract_image_from_shape(shape, slide_num, image_counter, output_dir, "fill")
                    if image_info:
                        images_info.append(image_info)
                        image_counter += 1
            except:
                pass
        
        # PLACEHOLDER shapes that might contain pictures
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            try:
                if hasattr(shape, 'image'):
                    image_info = extract_image_from_shape(shape, slide_num, image_counter, output_dir, "placeholder")
                    if image_info:
                        images_info.append(image_info)
                        image_counter += 1
            except:
                pass
    
    return images_info

def main():
    # Define paths
    pptx_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_covers.pptx")
    output_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/images")
    
    # Create output directory
    output_dir.mkdir(exist_ok=True)
    
    # Check if file exists
    if not pptx_file.exists():
        print(f"Error: File not found: {pptx_file}")
        return
    
    print(f"Loading presentation from: {pptx_file}")
    print(f"Output directory: {output_dir}\n")
    
    # Load presentation
    prs = Presentation(str(pptx_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides in presentation\n")
    
    # Extract images from all slides
    all_extractions = {}
    total_images = 0
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {idx}/{total_slides}...")
        
        try:
            images_info = extract_all_images_from_slide(slide, idx, output_dir)
            
            if images_info:
                all_extractions[f"slide_{idx:02d}"] = {
                    "slide_number": idx,
                    "image_count": len(images_info),
                    "images": images_info
                }
                total_images += len(images_info)
            else:
                print(f"    - No images found")
                
        except Exception as e:
            print(f"  ✗ Error processing slide {idx}: {e}")
    
    # Create index file with all metadata
    image_index = {
        "source_file": str(pptx_file.name),
        "extraction_timestamp": "2025-11-04",
        "total_slides": total_slides,
        "slides_with_images": len(all_extractions),
        "total_images_extracted": total_images,
        "extractions": all_extractions
    }
    
    index_file = output_dir / "image_index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(image_index, f, indent=2, ensure_ascii=False)
    
    # Print summary
    print(f"\n{'='*60}")
    print(f"✓ Image extraction complete!")
    print(f"{'='*60}")
    print(f"  Source file: {pptx_file.name}")
    print(f"  Total slides: {total_slides}")
    print(f"  Slides with images: {len(all_extractions)}")
    print(f"  Total images extracted: {total_images}")
    print(f"  Output directory: {output_dir}")
    print(f"  Index file: {index_file}")
    print(f"{'='*60}\n")

if __name__ == "__main__":
    main()



