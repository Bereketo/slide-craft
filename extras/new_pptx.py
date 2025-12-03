import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
data={
  "deck": {
    "title": "Q3 Business Review – AI & Compliance",
    "author": "Mobifly",
    "company": "PwC (Example)",
    "confidentiality": "Confidential",
    "date": "2025-09-18",
    "locale": "en-US",
    "rtl": False,
    "slide_size": "16x9",
    "slide_numbering": True,
    "header": "PwC | Mobifly – AI Automation",
    "footer_left": "© 2025 Mobifly",
    "footer_right": "Confidential",
    "watermark": { "text": "CONFIDENTIAL", "opacity": 0.08, "angle": -30 },
    "security": { "allow_remote_assets": False, "max_image_bytes": 5242880 }
  },

  "tokens": {
    "color": {
      "bg": "#0F172A",
      "surface": "#111827",
      "primary": "#2563EB",
      "accent": "#10B981",
      "danger": "#EF4444",
      "text": "#E5E7EB",
      "muted": "#9CA3AF",
      "tableHeader": "#1F2937",
      "zebra": "#0B1220",
      "border": "#374151"
    },
    "font": {
      "title_family": "Calibri",
      "body_family": "Calibri",
      "code_family": "Consolas",
      "fallbacks": ["Arial", "Noto Sans"],
      "title_size": 44,
      "h2_size": 28,
      "body_size": 18,
      "min_body_size": 14,
      "line_spacing": 1.2
    },
    "spacing": { "margin": 24, "gutter": 12 },
    "grid": { "columns": 12, "unit": "px" }
  },

  "defaults": {
    "title_style": { "bold": True, "color": "text", "align": "left" },
    "body_style": { "color": "text", "align": "left" },
    "table_style": {
      "header_fill": "#1F2937",
      "header_color": "#E5E7EB",
      "row_zebra": "#0B1220",
      "border_color": "#374151"
    },
    "chart_style": {
      "palette": ["#2563EB", "#10B981", "#F59E0B", "#EF4444"],
      "data_labels": True,
      "legend": "right"
    }
  },

  "slides": [
    {
      "title": "Cover",
      "background": { "type": "solid", "color": "bg" },
      "components": [
        {
          "type": "text",
          "text_type": "title",
          "value": "Q3 Business Review",
          "grid": { "col": 1, "span": 12, "row_h": 120, "y": 120 },
          "style": { "align": "left" },
          "alt": "Deck title"
        },
        {
          "type": "text",
          "text_type": "body",
          "value": "AI & Compliance Highlights",
          "grid": { "col": 1, "span": 12, "row_h": 60, "y": 220 },
          "style": { "font_size": 22, "color": "#9CA3AF" }
        },
        {
          "type": "image",
          "src": "assets/pwc_lockup.png",
          "grid": { "col": 1, "span": 4, "row_h": 80, "y": 24 },
          "object_fit": "contain",
          "alt": "PwC logo"
        }
      ],
      "notes": "Tailor talking points to client industry; call out measurable outcomes."
    },

    {
      "title": "Executive Summary",
      "background": { "type": "solid", "color": "bg" },
      "components": [
        {
          "type": "text",
          "text_type": "h2",
          "value": "Key Outcomes",
          "grid": { "col": 1, "span": 12, "row_h": 60, "y": 24 }
        },
        {
          "type": "richtext",
          "runs": [
            {"text": "• 12% QoQ revenue growth; ", "bold": True},
            {"text": "AI-driven invoice validation reduced exceptions by 37%.\n"},
            {"text": "• 7 new logos in BFSI & GCC; ", "bold": True},
            {"text": "pipeline strengthened in KSA.\n"},
            {"text": "• Compliance latency cut from 3 days to ", "bold": False},
            {"text": "6 hours", "bold": True},
            {"text": ".", "bold": False}
          ],
          "grid": { "col": 1, "span": 12, "row_h": 200, "y": 84 },
          "fit": "wrap",
          "style": { "font_size": 20 }
        }
      ]
    },

    {
      "title": "Key Metrics",
      "background": { "type": "solid", "color": "bg" },
      "components": [
        {
          "type": "text",
          "text_type": "h2",
          "value": "KPIs",
          "grid": { "col": 1, "span": 6, "row_h": 48, "y": 24 }
        },
        {
          "type": "table",
          "grid": { "col": 1, "span": 6, "row_h": 320, "y": 80 },
          "fit": "paginate",
          "content": {
            "columns": ["Metric", "Q3", "QoQ"],
            "rows": [
              ["Revenue", "$2.4M", "+12%"],
              ["New Logos", "7", "+2"],
              ["Churn", "1.2%", "-0.3 pp"],
              ["Invoices Auto-cleared", "78%", "+11 pp"],
              ["Avg Ticket Cycle Time", "6h", "-18h"]
            ]
          },
          "table_options": {
            "column_widths": [240, 140, 120],
            "header_row": True,
            "first_col_bold": True
          },
          "alt": "Table of key performance indicators"
        },
        {
          "type": "chart",
          "grid": { "col": 7, "span": 6, "row_h": 360, "y": 80 },
          "fit": "wrap",
          "content": {
            "categories": ["India", "UAE", "KSA"],
            "series": [
              {"name": "Leads", "values": [42, 18, 11]},
              {"name": "Qualified", "values": [21, 9, 5]}
            ]
          },
          "chart_options": { "chartType": "column", "legend": "right", "data_labels": True },
          "alt": "Column chart of leads and qualified by region"
        }
      ]
    },

    {
      "title": "Architecture Snapshot",
      "background": { "type": "solid", "color": "bg" },
      "components": [
        {
          "type": "text",
          "text_type": "h2",
          "value": "Invoice Validation Platform (High Level)",
          "grid": { "col": 1, "span": 12, "row_h": 60, "y": 24 }
        },
        {
          "type": "image",
          "src": "assets/arch-diagram.png",
          "grid": { "col": 1, "span": 12, "row_h": 420, "y": 84 },
          "object_fit": "contain",
          "alt": "High-level architecture diagram"
        }
      ],
      "notes": "Ensure client-approved diagram; avoid exposing internal hostnames."
    },

    {
      "title": "Appendix: Sample Code",
      "background": { "type": "solid", "color": "bg" },
      "components": [
        {
          "type": "text",
          "text_type": "h2",
          "value": "Policy Check (Python)",
          "grid": { "col": 1, "span": 12, "row_h": 48, "y": 24 }
        },
        {
          "type": "code",
          "value": "def check_threshold(invoice):\n    return invoice.amount <= 500000 and invoice.pan_verified\n",
          "code_options": { "language": "python", "render_as": "image" },
          "grid": { "col": 1, "span": 12, "row_h": 300, "y": 84 },
          "alt": "Sample Python code block rendering as image"
        }
      ]
    }
  ]
}

def create_ppt_from_json(data: dict, output_file="output.pptx"):
    # Initialize presentation
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)  # 16:9

    # Grid system
    cols = data["tokens"]["grid"]["columns"]
    margin = data["tokens"]["spacing"]["margin"]
    gutter = data["tokens"]["spacing"]["gutter"]

    col_width = (prs.slide_width - 2 * Inches(margin/72)) / cols  # px → pt → inch

    def grid_to_position(grid):
        """Convert grid spec to pptx position + size"""
        x = Inches(margin/72) + col_width * (grid["col"] - 1)
        w = col_width * grid["span"]
        y = Inches(grid["y"]/72)
        h = Inches(grid["row_h"]/72)
        return x, y, w, h

    # Loop over slides
    for slide_json in data["slides"]:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Background
        bg_color = data["tokens"]["color"].get(slide_json["background"]["color"], "#FFFFFF")
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(bg_color.strip("#"))

        # Components
        for comp in slide_json["components"]:
            x, y, w, h = grid_to_position(comp["grid"])

            if comp["type"] == "text":
                box = slide.shapes.add_textbox(x, y, w, h)
                tf = box.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = comp["value"]
                if "style" in comp:
                    style = comp["style"]
                    if "font_size" in style:
                        run.font.size = Pt(style["font_size"])
                    if "color" in style:
                        run.font.color.rgb = RGBColor.from_string(
                            style["color"].strip("#") if style["color"].startswith("#")
                            else data["tokens"]["color"].get(style["color"], "#E5E7EB").strip("#")
                        )
                tf.word_wrap = True

            elif comp["type"] == "image":
                try:
                    slide.shapes.add_picture(comp["src"], x, y, width=w, height=h)
                except:
                    print(f"Image not found :: {comp['src']}")
            elif comp["type"] == "table":
                rows = len(comp["content"]["rows"]) + 1
                cols_t = len(comp["content"]["columns"])
                table = slide.shapes.add_table(rows, cols_t, x, y, w, h).table

                # header
                for j, col_name in enumerate(comp["content"]["columns"]):
                    cell = table.cell(0, j)
                    cell.text = col_name
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(
                        data["defaults"]["table_style"]["header_fill"].strip("#")
                    )
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = RGBColor.from_string(
                            data["defaults"]["table_style"]["header_color"].strip("#")
                        )

                # rows
                for i, row in enumerate(comp["content"]["rows"]):
                    for j, val in enumerate(row):
                        table.cell(i + 1, j).text = str(val)

            # TODO: implement richtext, charts, code rendering

    prs.save(output_file)
    print(f"PPT generated: {output_file}")
create_ppt_from_json(data,)