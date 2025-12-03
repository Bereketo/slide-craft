#!/usr/bin/env python3
"""
Extract all images from PwC cover slides and save them.
This will extract the actual photos, logos, and graphics from each cover.
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def extract_images_from_slide(slide, slide_num, output_dir):
    """Extract all images from a slide and save them"""
    
    images_info = []
    image_count = 0
    
    for shape in slide.shapes:
        # Check if shape is a picture
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            
            try:
                # Get the image
                image = shape.image
                image_bytes = image.blob
                
                # Determine file extension from content type
                ext_map = {
                    'image/jpeg': 'jpg',
                    'image/jpg': 'jpg',
                    'image/png': 'png',
                    'image/gif': 'gif',
                    'image/bmp': 'bmp',
                    'image/tiff': 'tiff'
                }
                ext = ext_map.get(image.content_type, 'png')
                
                # Save the image
                image_filename = f"cover_{slide_num:02d}_image_{image_count}.{ext}"
                image_path = output_dir / image_filename
                
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                # Get image dimensions and position
                left_px = int(shape.left / 9525)
                top_px = int(shape.top / 9525)
                width_px = int(shape.width / 9525)
                height_px = int(shape.height / 9525)
                
                # Get actual image dimensions
                try:
                    img = Image.open(io.BytesIO(image_bytes))
                    actual_width, actual_height = img.size
                except:
                    actual_width, actual_height = width_px, height_px
                
                images_info.append({
                    "filename": image_filename,
                    "position": {
                        "x": left_px,
                        "y": top_px,
                        "w": width_px,
                        "h": height_px
                    },
                    "actual_dimensions": {
                        "width": actual_width,
                        "height": actual_height
                    },
                    "content_type": image.content_type
                })
                
                print(f"    ✓ Extracted: {image_filename} ({actual_width}x{actual_height})")
                
            except Exception as e:
                print(f"    ✗ Error extracting image {image_count}: {e}")
    
    return images_info

def update_cover_json_with_images(cover_json_path, images_info, images_relative_path):
    """Update the cover JSON file to reference the actual extracted images"""
    
    with open(cover_json_path, 'r', encoding='utf-8') as f:
        cover_data = json.load(f)
    
    # Update image components with actual image paths
    components = cover_data['slide']['components']
    image_index = 0
    
    for component in components:
        if component['type'] == 'image' and image_index < len(images_info):
            image_info = images_info[image_index]
            component['src'] = f"{images_relative_path}/{image_info['filename']}"
            
            # Update position to match actual extracted position
            if 'box' in component:
                component['box'] = {
                    "x": image_info['position']['x'],
                    "y": image_info['position']['y'],
                    "w": image_info['position']['w'],
                    "h": image_info['position']['h'],
                    "unit": "px"
                }
            
            # Add metadata
            component['image_metadata'] = {
                "actual_width": image_info['actual_dimensions']['width'],
                "actual_height": image_info['actual_dimensions']['height'],
                "content_type": image_info['content_type']
            }
            
            image_index += 1
    
    # Save updated JSON
    with open(cover_json_path, 'w', encoding='utf-8') as f:
        json.dump(cover_data, f, indent=2, ensure_ascii=False)

def main():
    # Paths
    covers_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_covers.pptx")
    covers_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/covers")
    images_dir = covers_dir / "images"
    
    # Create images directory
    images_dir.mkdir(exist_ok=True)
    
    if not covers_file.exists():
        print(f"Error: File not found: {covers_file}")
        return
    
    print(f"Loading presentation from: {covers_file}")
    prs = Presentation(str(covers_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides in presentation")
    print(f"Extracting images to: {images_dir}\n")
    
    # Extract images from each slide
    all_extractions = {}
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {idx}/{total_slides}...")
        
        try:
            images_info = extract_images_from_slide(slide, idx, images_dir)
            
            if images_info:
                all_extractions[idx] = images_info
                
                # Update the corresponding JSON file
                json_file = covers_dir / f"cover_{idx:02d}_title_subtitle.json"
                if json_file.exists():
                    update_cover_json_with_images(json_file, images_info, "images")
                    print(f"    ✓ Updated JSON: {json_file.name}")
            else:
                print(f"    - No images found")
                
        except Exception as e:
            print(f"  ✗ Error processing slide {idx}: {e}")
    
    # Create image index
    image_index = {
        "total_covers": total_slides,
        "covers_with_images": len(all_extractions),
        "total_images_extracted": sum(len(imgs) for imgs in all_extractions.values()),
        "extractions": {
            f"cover_{idx:02d}": {
                "slide_number": idx,
                "image_count": len(images),
                "images": images
            }
            for idx, images in all_extractions.items()
        }
    }
    
    index_file = images_dir / "image_index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(image_index, f, indent=2, ensure_ascii=False)
    
    print(f"\n✓ Image extraction complete!")
    print(f"  Total covers: {total_slides}")
    print(f"  Covers with images: {len(all_extractions)}")
    print(f"  Total images extracted: {image_index['total_images_extracted']}")
    print(f"  Images directory: {images_dir}")
    print(f"  Image index: {index_file}")

if __name__ == "__main__":
    main()







