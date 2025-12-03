#!/usr/bin/env python3
"""
Extract PwC cover slides from the template PPTX and convert to JSON format.
Each cover will be saved as a separate JSON file with metadata about its use case.
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_px(emu_value):
    """Convert EMU to pixels (96 DPI)"""
    return int(emu_value / 9525)

def rgb_to_hex(rgb_color):
    """Convert RGB color to hex"""
    try:
        return f"#{rgb_color.r:02X}{rgb_color.g:02X}{rgb_color.b:02X}"
    except:
        return "#000000"

def get_fill_color(shape):
    """Extract fill color from shape"""
    try:
        if shape.fill.type == 1:  # SOLID
            return rgb_to_hex(shape.fill.fore_color.rgb)
    except:
        pass
    return None

def extract_text_from_shape(shape):
    """Extract text and styling from a shape"""
    if not shape.has_text_frame:
        return None
    
    text_frame = shape.text_frame
    if not text_frame.text.strip():
        return None
    
    components = []
    
    for paragraph in text_frame.paragraphs:
        if not paragraph.text.strip():
            continue
            
        # Check if it's simple text or needs richtext
        runs = []
        for run in paragraph.runs:
            if not run.text:
                continue
                
            # Extract font size - try multiple ways
            font_size = 18  # default
            try:
                if run.font.size:
                    font_size = int(run.font.size.pt)
            except:
                pass
            
            run_info = {
                "text": run.text,
                "bold": run.font.bold if run.font.bold is not None else False,
                "font_size": font_size,
            }
            
            # Try to get font family
            try:
                if run.font.name:
                    run_info["font_family"] = run.font.name
            except:
                try:
                    # Try alternate method
                    font_name = str(run.font._element.latin.typeface) if hasattr(run.font._element, 'latin') else None
                    if font_name:
                        run_info["font_family"] = font_name
                except:
                    pass
            
            # Try to get color - handle multiple color types
            try:
                color_obj = run.font.color
                if color_obj.type == 1:  # RGB (MSO_COLOR_TYPE.RGB)
                    run_info["color"] = rgb_to_hex(color_obj.rgb)
                elif color_obj.type == 2:  # SCHEME (MSO_COLOR_TYPE.SCHEME)
                    # Try to get the actual RGB value from scheme color
                    try:
                        run_info["color"] = rgb_to_hex(color_obj.rgb)
                    except:
                        pass
                elif hasattr(color_obj, 'rgb'):
                    # Fallback: try to get RGB directly
                    run_info["color"] = rgb_to_hex(color_obj.rgb)
            except:
                pass
                
            runs.append(run_info)
        
        if runs:
            components.append({
                "runs": runs,
                "alignment": str(paragraph.alignment) if paragraph.alignment else "LEFT"
            })
    
    return components if components else None

def categorize_slide(slide_num, slide, text_content):
    """Categorize the slide based on its content and layout"""
    
    # Common categories based on PwC cover templates
    categories = {
        "title_only": "Simple title cover",
        "title_subtitle": "Title with subtitle",
        "title_date": "Title with date",
        "title_author": "Title with author/presenter",
        "title_image": "Title with image/graphic",
        "section_divider": "Section divider",
        "agenda": "Agenda/Contents",
        "full_image": "Full image background",
        "branded": "PwC branded template",
        "minimal": "Minimal design",
        "executive": "Executive summary",
        "data_focus": "Data-focused layout"
    }
    
    # Simple categorization logic
    text_lower = text_content.lower()
    
    if "agenda" in text_lower or "content" in text_lower:
        return "agenda", "Agenda or contents page"
    elif "section" in text_lower or len(text_content) < 30:
        return "section_divider", "Section divider or chapter title"
    elif "executive" in text_lower:
        return "executive", "Executive summary or overview"
    else:
        # Count shapes to determine complexity
        shape_count = len([s for s in slide.shapes if s.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER])
        if shape_count > 5:
            return "title_image", "Title with decorative elements"
        else:
            return "title_subtitle", "Standard title with subtitle"

def extract_slide_to_json(slide, slide_num):
    """Extract a single slide and convert to our JSON schema format"""
    
    components = []
    all_text = ""
    
    # Extract all shapes
    for shape in slide.shapes:
        # Skip placeholders that are empty
        if shape.is_placeholder and not shape.has_text_frame:
            continue
        
        # Get position and size
        left_px = emu_to_px(shape.left)
        top_px = emu_to_px(shape.top)
        width_px = emu_to_px(shape.width)
        height_px = emu_to_px(shape.height)
        
        # Extract text components
        if shape.has_text_frame:
            text_components = extract_text_from_shape(shape)
            if text_components:
                all_text += " " + shape.text
                
                for text_comp in text_components:
                    # Determine if it's title, h2, or body based on font size
                    avg_size = sum(r.get("font_size", 18) for r in text_comp["runs"]) / len(text_comp["runs"])
                    
                    if avg_size >= 36:
                        text_type = "title"
                    elif avg_size >= 24:
                        text_type = "h2"
                    else:
                        text_type = "body"
                    
                    component = {
                        "type": "richtext",
                        "runs": text_comp["runs"],
                        "box": {
                            "x": left_px,
                            "y": top_px,
                            "w": width_px,
                            "h": height_px,
                            "unit": "px"
                        },
                        "style": {
                            "align": text_comp["alignment"].lower() if text_comp["alignment"] else "left"
                        }
                    }
                    
                    components.append(component)
        
        # Extract image components
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            component = {
                "type": "image",
                "src": "placeholder_image.png",  # Placeholder - would need to extract actual image
                "alt": f"Cover image {slide_num}",
                "box": {
                    "x": left_px,
                    "y": top_px,
                    "w": width_px,
                    "h": height_px,
                    "unit": "px"
                },
                "object_fit": "cover"
            }
            components.append(component)
        
        # Extract shape components
        elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
            fill_color = get_fill_color(shape)
            if fill_color:
                component = {
                    "type": "shape",
                    "shape_type": "rectangle",
                    "box": {
                        "x": left_px,
                        "y": top_px,
                        "w": width_px,
                        "h": height_px,
                        "unit": "px"
                    },
                    "style": {
                        "fill": fill_color
                    }
                }
                components.append(component)
    
    # Categorize the slide
    category, description = categorize_slide(slide_num, slide, all_text)
    
    return {
        "metadata": {
            "slide_number": slide_num,
            "category": category,
            "description": description,
            "extracted_text": all_text.strip()[:200]  # First 200 chars for reference
        },
        "slide": {
            "title": f"PwC Cover {slide_num}",
            "background": {"type": "solid", "color": "#FFFFFF"},
            "components": components
        }
    }

def main():
    # Path to the PwC covers file
    covers_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_covers.pptx")
    
    if not covers_file.exists():
        print(f"Error: File not found: {covers_file}")
        return
    
    # Create output directory
    output_dir = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/covers")
    output_dir.mkdir(exist_ok=True)
    
    print(f"Loading presentation from: {covers_file}")
    prs = Presentation(str(covers_file))
    
    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides in presentation")
    
    # Extract each slide
    all_covers = []
    category_counts = {}
    
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {idx}/{total_slides}...")
        
        try:
            cover_json = extract_slide_to_json(slide, idx)
            
            # Save individual cover
            output_file = output_dir / f"cover_{idx:02d}_{cover_json['metadata']['category']}.json"
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(cover_json, f, indent=2, ensure_ascii=False)
            
            all_covers.append(cover_json)
            
            # Track categories
            category = cover_json['metadata']['category']
            category_counts[category] = category_counts.get(category, 0) + 1
            
            print(f"  ✓ Saved: {output_file.name}")
            
        except Exception as e:
            print(f"  ✗ Error processing slide {idx}: {e}")
    
    # Create index file
    index = {
        "total_covers": len(all_covers),
        "categories": category_counts,
        "covers": [
            {
                "id": cover["metadata"]["slide_number"],
                "file": f"cover_{cover['metadata']['slide_number']:02d}_{cover['metadata']['category']}.json",
                "category": cover["metadata"]["category"],
                "description": cover["metadata"]["description"]
            }
            for cover in all_covers
        ]
    }
    
    index_file = output_dir / "index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(index, f, indent=2, ensure_ascii=False)
    
    print(f"\n✓ Extraction complete!")
    print(f"  Total covers extracted: {len(all_covers)}")
    print(f"  Output directory: {output_dir}")
    print(f"  Index file: {index_file}")
    print(f"\nCategory breakdown:")
    for category, count in sorted(category_counts.items()):
        print(f"  - {category}: {count} covers")

if __name__ == "__main__":
    main()





