#!/usr/bin/env python3
"""
Debug script to check font properties in slide 2 of the content elements file.
"""

from pptx import Presentation
from pathlib import Path

def inspect_slide_2():
    pptx_file = Path("/Users/ahmshalan/Projects/mobifly/ppt-generate/branding/PwC_ppt_graphic_elements_level1.pptx")
    prs = Presentation(str(pptx_file))
    
    # Get slide 2 (index 1)
    slide = prs.slides[1]
    
    print("=== SLIDE 2 INSPECTION ===\n")
    
    for idx, shape in enumerate(slide.shapes):
        print(f"\n--- Shape {idx} ---")
        print(f"Shape Type: {shape.shape_type}")
        print(f"Has Text Frame: {shape.has_text_frame}")
        
        if shape.has_text_frame:
            text_frame = shape.text_frame
            print(f"Text: '{text_frame.text}'")
            
            for p_idx, paragraph in enumerate(text_frame.paragraphs):
                print(f"\n  Paragraph {p_idx}: '{paragraph.text}'")
                
                for r_idx, run in enumerate(paragraph.runs):
                    print(f"\n    Run {r_idx}:")
                    print(f"      Text: '{run.text}'")
                    print(f"      Font.name: {run.font.name}")
                    print(f"      Font.size: {run.font.size}")
                    if run.font.size:
                        print(f"      Font.size.pt: {run.font.size.pt}")
                        print(f"      Font.size.emu: {run.font.size}")
                    print(f"      Font.bold: {run.font.bold}")
                    print(f"      Font.color.type: {run.font.color.type}")
                    
                    # Try to get color RGB
                    try:
                        rgb = run.font.color.rgb
                        print(f"      Font.color.rgb: #{rgb.r:02X}{rgb.g:02X}{rgb.b:02X}")
                    except:
                        print(f"      Font.color.rgb: Could not extract")
                    
                    # Check for font family in XML
                    try:
                        if hasattr(run.font._element, 'rPr'):
                            rPr = run.font._element.rPr
                            if rPr is not None:
                                print(f"      rPr exists: {rPr}")
                                # Try to find latin font
                                latin = rPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                                if latin is not None:
                                    print(f"      Latin font: {latin.get('typeface')}")
                    except Exception as e:
                        print(f"      XML inspection error: {e}")

if __name__ == "__main__":
    inspect_slide_2()




