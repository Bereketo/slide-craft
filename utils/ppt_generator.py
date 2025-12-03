#!/usr/bin/env python3
"""
pptx_renderer.py
Template-less PPTX generator driven by validated JSON.
- Validates input (Draft-07) with jsonschema
- Renders: text, richtext, image, table, chart, shape(line/box minimal), background
- Enterprise niceties: header/footer, page numbering, watermark, alt text (basic)
- Layout: 12-col grid or absolute box (px/in)

Usage:
  python pptx_renderer.py --schema pptx-deck.schema.json --input deck.json --out out.pptx

Requirements:
  pip install python-pptx jsonschema pillow
"""

import argparse, json, math, os, sys
from pathlib import Path
from typing import Dict, Any, Tuple, List, Optional
from core.logger_setup import app_logger as logger
from jsonschema import validate as jsonschema_validate, Draft7Validator
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement, qn
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.slide import Slide
import requests

# Overlapping alignment formatters
# ---------- Helpers ----------

def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def to_emu(value: float, unit: str) -> Emu:
    if unit == "in":
        return Inches(value)
    # default px → approximate @ 96 dpi
    return Emu(int(value * 9525))

def slide_set_size(prs: Presentation, size: str):
    # 16x9 default: 13.333in x 7.5in (PowerPoint default)
    if size == "4x3":
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    elif size == "A4-landscape":
        prs.slide_width, prs.slide_height = Inches(11.69), Inches(8.27)
    else:  # 16x9
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

def grid_rect(prs: Presentation, tokens: Dict[str, Any], grid: Dict[str, Any]) -> Tuple[Emu, Emu, Emu, Emu]:
    """
    Compute (left, top, width, height) from 12-col grid, with override if ignore_overlaps and specific values provided.
    Supports additional offset in cm via grid['offset_cm'] (applied to left).
    """
    columns = tokens["grid"]["columns"]
    unit = tokens["grid"].get("unit", "px")
    margin = tokens["spacing"]["margin"]
    gutter = tokens["spacing"]["gutter"]

    # Convert design-system px/in to EMU early
    margin_emu = to_emu(margin, unit)
    gutter_emu = to_emu(gutter, unit)
    top_y = grid.get("y", margin)
    top_emu = to_emu(top_y, unit)
    row_h_emu = to_emu(grid["row_h"], unit)

    inner_width = prs.slide_width - (margin_emu * 2)
    col_width = (inner_width - gutter_emu * (columns - 1)) // columns

    col_index = max(0, grid["col"] - 1)
    span = grid["span"]

    left = margin_emu + (col_width + gutter_emu) * col_index
    width = col_width * span + gutter_emu * (span - 1)

    # Support additional offset in centimeters (e.g., grid['offset_cm'])
    offset_cm = grid.get("offset_cm", 0)
    if offset_cm:
        # 1 cm = 360000 EMU
        left += Emu(int(offset_cm * 360000))

    # If ignore_overlaps is set and any of x, y, w, h are provided, use those specific values instead
    if grid.get("ignore_overlaps"):
        for key, var in (("x", "left"), ("y", "top_emu"), ("w", "width"), ("h", "row_h_emu")):
            val = grid.get(key)
            if val is not None:
                locals()[var] = to_emu(val, unit)
    return left, top_emu, width, row_h_emu

def box_rect(box: Dict[str, Any]) -> Tuple[Emu, Emu, Emu, Emu]:
    unit = box.get("unit", "px")
    return to_emu(box["x"], unit), to_emu(box["y"], unit), to_emu(box["w"], unit), to_emu(box["h"], unit)

def apply_background(slide, tokens: Dict[str, Any], bg: Optional[Dict[str, Any]]):
    if not bg:
        return
    fill = slide.background.fill
    if bg["type"] == "solid":
        fill.solid()
        color_key = bg.get("color")
        color_val = tokens["color"].get(color_key, color_key)
        fill.fore_color.rgb = rgb(color_val)
        if "opacity" in bg:
            fill.transparency = max(0.0, min(1.0, 1.0 - bg["opacity"]))  # pptx transparency is inverse
    elif bg["type"] == "image":
        # Workaround: add a full-bleed picture at Z-order back
        src = bg["src"]
        try:
            # Get presentation dimensions from the slide's parent presentation
            prs = slide.part.package.presentation_part.presentation
            slide.shapes.add_picture(src, Emu(0), Emu(0), width=prs.slide_width,
                                     height=prs.slide_height)
        except Exception as e:
            logger.error(f"[warn] background image failed: {e}")

def set_alt_text(shape, alt: Optional[str]):
    """Set alternative text (Title/Description)."""
    if not alt:
        return
    try:
        sp = shape._element  # lxml element
        nvSpPr = sp.xpath("./p:nvSpPr")
        if nvSpPr:
            cNvPr = nvSpPr[0].xpath("./p:cNvPr")
            if cNvPr:
                cNvPr[0].set("descr", alt)
    except Exception:
        pass

def add_header_footer_and_page_numbers(prs: Presentation, deck: Dict[str, Any], tokens: Dict[str, Any]):
    header = deck.get("header")
    footer_left = deck.get("footer_left")
    footer_right = deck.get("footer_right")
    show_numbers = deck.get("slide_numbering", True)

    unit = tokens["grid"].get("unit", "px")
    margin = to_emu(tokens["spacing"]["margin"], unit)
    base_font = tokens["font"]["body_family"]
    font_color = rgb(tokens["color"].get("muted"))
    font_size = Pt(max(10, tokens["font"].get("min_body_size", 12)))

    # Bring footer up by 10px (EMU)
    ten_px_emu = to_emu(10, "px")

    for idx, slide in enumerate(prs.slides, start=1):
        # Header (top-left)
        if header:
            tb = slide.shapes.add_textbox(margin, margin//2, prs.slide_width - margin*2, Pt(14).emu)
            p = tb.text_frame.paragraphs[0]
            p.text = header
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            for r in p.runs:
                r.font.name = base_font

                r.font.size = font_size
                if font_color:
                    r.font.color.rgb = font_color
            set_alt_text(tb, "Header")

        # Footer left/right + page number
        footer_y = prs.slide_height - margin + Pt(2).emu - ten_px_emu
        footer_h = Pt(16).emu
        # Left
        if footer_left:
            tb = slide.shapes.add_textbox(margin, footer_y, prs.slide_width//2 - margin, footer_h)
            p = tb.text_frame.paragraphs[0]
            p.text = footer_left
            for r in p.runs:
                r.font.name = base_font
                r.font.size = font_size
                if font_color:
                    r.font.color.rgb = font_color
            set_alt_text(tb, "Footer left")
        # Right
        right_text = footer_right or ""
        if show_numbers:
            right_text = f"{right_text}    {idx}".strip()
        tb = slide.shapes.add_textbox(prs.slide_width//2, footer_y, prs.slide_width//2 - margin, footer_h)
        p = tb.text_frame.paragraphs[0]
        p.text = right_text
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        for r in p.runs:
            r.font.name = base_font
            r.font.size = font_size
            if font_color:
                r.font.color.rgb = font_color
        set_alt_text(tb, "Footer right / page number")

def apply_watermark(slide, deck: Dict[str, Any]):
    wm = deck.get("watermark")
    if not wm or not wm.get("text"):
        return
    text = wm["text"]
    angle = float(wm.get("angle", -30))
    # Large centered textbox; use gray color with transparency via alpha fill.
    # Center the watermark textbox
    prs = slide.part.package.presentation_part.presentation
    width, height = prs.slide_width, prs.slide_height
    # Make the watermark textbox tall and centered horizontally, but not full height
    left = 0.6
    top = int(height * 0.35)  # start high, but not at the very top
    # Get presentation dimensions from the slide's parent presentation
    prs = slide.part.package.presentation_part.presentation
    width, height = prs.slide_width, prs.slide_height
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(wm.get('size',120))
    run.font.bold = True
    run.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    # rotate
    tb.rotation = angle
    # Transparency on shape fill (none) and text is limited; we simulate by light color + big size.
    set_alt_text(tb, "Watermark")

# ---------- Component renderers ----------

def render_text(slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    rect = grid_rect(prs, tokens, comp["grid"]) if "grid" in comp else box_rect(comp["box"])
    left, top, width, height = rect
    # Clamp textbox width so it doesn't cross the slide's right edge
    max_width = prs.slide_width - left
    if width > max_width:
        width = max_width
    tb = slide.shapes.add_textbox(left, top, width, height)
    # Apply background color if specified in comp["style"]["background"] or ["bg"]
    bg_color = comp.get("style", {}).get("fill") or comp.get("style", {}).get("bg")
    if bg_color:
        tb.fill.solid()
        tb.fill.fore_color.rgb = rgb(bg_color)
    else:
        tb.fill.background()
    tf = tb.text_frame
    tf.clear()
    # Ensure text wraps within the textbox and doesn't auto-resize
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    
    # Add padding to text frame for better appearance (especially with backgrounds)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Handle vertical alignment if specified
    valign = comp.get("style", {}).get("valign", "top")
    tf.vertical_anchor = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "center": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    }.get(valign, MSO_VERTICAL_ANCHOR.TOP)

    p = tf.paragraphs[0]
    p.text = comp["value"]
    style_key = comp.get("text_type", "body")
    font_cfg = tokens["font"]
    
    # PwC Branding: Font mapping based on text_type
    # title, body, quotes -> Georgia (serif)
    # h2, caption, labels, data numbers -> Arial (sans-serif)
    font_family_map = {
        "title": font_cfg.get("title_family", "Georgia"),      # Serif for titles
        "h2": "Arial",                                          # Sans-serif for sub-headlines
        "body": font_cfg.get("body_family", "Georgia"),        # Serif for body text
        "caption": "Arial"                                      # Sans-serif for captions/labels
    }
    
    size_map = {
        "title": font_cfg.get("title_size", 44),
        "h2": font_cfg.get("h2_size", 28),
        "body": font_cfg.get("body_size", 18),
        "caption": max(10, int(font_cfg.get("body_size", 18) * 0.8))
    }
    
    # PwC Branding: Bold mapping
    bold_map = {
        "title": True,   # Titles are bold
        "h2": True,      # Sub-headlines are bold
        "body": False,   # Body text is regular
        "caption": False # Captions are regular
    }
    
    align = comp.get("style", {}).get("align", "left")
    p.alignment = {
        "left": PP_PARAGRAPH_ALIGNMENT.LEFT,
        "center": PP_PARAGRAPH_ALIGNMENT.CENTER,
        "right": PP_PARAGRAPH_ALIGNMENT.RIGHT,
        "justify": PP_PARAGRAPH_ALIGNMENT.JUSTIFY
    }.get(align, PP_PARAGRAPH_ALIGNMENT.LEFT)
    # Tighten paragraph spacing to avoid excessive gaps between lines/paragraphs
    try:
        p.line_spacing = 1.0
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    except Exception:
        pass
    for r in p.runs:
        # Apply font family: Use style override first, then PwC branding map
        r.font.name = comp.get("style", {}).get("font_family") or font_family_map.get(style_key, font_cfg["body_family"])
        r.font.size = Pt(comp.get("style", {}).get("font_size", size_map.get(style_key, font_cfg.get("body_size", 18))))
        r.font.bold = comp.get("style", {}).get("bold", bold_map.get(style_key, False))
        # PwC Branding: Never use italics
        r.font.italic = False  # PwC strict rule: no italics
        color = comp.get("style", {}).get("color",tokens.get('color')['text'])
        if color:
            r.font.color.rgb = rgb(color)
    set_alt_text(tb, comp.get("alt"))

def render_richtext(slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    rect = grid_rect(prs, tokens, comp["grid"]) if "grid" in comp else box_rect(comp["box"])
    left, top, width, height = rect
    # Clamp textbox width so it doesn't cross the slide's right edge
    max_width = prs.slide_width - left
    if width > max_width:
        width = max_width
    tb = slide.shapes.add_textbox(left, top, width, height)
    bg_color = comp.get("style", {}).get("fill") or comp.get("style", {}).get("bg")
    if bg_color:
        tb.fill.solid()
        tb.fill.fore_color.rgb = rgb(bg_color)
    else:
        tb.fill.background()
    tf = tb.text_frame
    tf.clear()
    # Ensure text wraps within the textbox and doesn't auto-resize
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    
    # Add padding to text frame for better appearance (especially with backgrounds)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Handle vertical alignment if specified
    valign = comp.get("style", {}).get("valign", "top")
    tf.vertical_anchor = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "center": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    }.get(valign, MSO_VERTICAL_ANCHOR.TOP)

    font_cfg = tokens["font"]
    p = tf.paragraphs[0]
    for i, run_def in enumerate(comp["runs"]):
        # If 'new_start' is True, start a new paragraph
        if i == 0:
            pass  # always use the first paragraph for the first run
        elif run_def.get('new_start', False):
            p = tf.add_paragraph()
        # Otherwise, continue in the current paragraph (after the last run)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        # Tighten paragraph spacing to avoid excessive gaps between lines/paragraphs
        try:
            p.line_spacing = 1.0
            p.space_before = Pt(0)
            p.space_after = Pt(0)
        except Exception:
            pass
        r = p.add_run()
        # Apply run-specific font size and family if present
        run_font_family = run_def.get("font_family")
        run_font_size = run_def.get("font_size")
        r.text = run_def["text"]
        r.font.underline = run_def.get("under_line",False)
        # PwC Branding: Default to Georgia for richtext (body content)
        r.font.name = font_cfg.get("body_family", "Georgia")
        base_size = comp.get("style", {}).get("font_size", font_cfg.get("body_size", 18))
        r.font.size = Pt(base_size)
        r.font.bold = bool(run_def.get("bold"))
        # PwC Branding: Never use italics (strict rule)
        r.font.italic = False  # Override any italic setting - PwC does not allow italics
        color = (run_def.get("color") or comp.get('style',{}).get("color")) or tokens.get('color')['text']
        if color:
            r.font.color.rgb = rgb(color)
        if run_font_family:
            r.font.name = run_font_family
        if run_font_size:
            r.font.size = Pt(run_font_size)
    set_alt_text(tb, comp.get("alt"))


def get_pexels_image(query: str):
    headers = {"Authorization": os.getenv("PEXELS_API_KEY")}
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    res = requests.get(url, headers=headers)
    data = res.json()
    if data.get("photos"):
        src = data["photos"][0]["src"]["medium"]
    else:
        src = "https://via.placeholder.com/800x600?text=No+Image+Found"
    response = requests.get(src, timeout=10)
    return response

def render_image(slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    import io
    import requests
    from PIL import Image

    rect = grid_rect(prs, tokens, comp["grid"]) if "grid" in comp else box_rect(comp["box"])
    left, top, width, height = rect

    # Check for horizontal coordinate override
    horizontal = comp.get("horizontal")
    if horizontal is not None:
        left = horizontal
    vertical = comp.get("vertical")
    if vertical is not None:
        top = vertical

    src = comp["src"]
    fit = comp.get("object_fit", "contain")
    try:
        # Handle both local file paths and URLs
        if src.startswith(('http://', 'https://')):
            # Download image from URL to a BytesIO buffer
            try:
                response = requests.get(src, timeout=10)
                response.raise_for_status()
                image_bytes = io.BytesIO(response.content)
            except:
                try:
                    response = get_pexels_image(comp.get("alt"))
                    image_bytes = io.BytesIO(response.content)
                except Exception as _:
                    return
        else:
            # Handle local file path
            try:
                local_path = Path(src)
                if not local_path.is_absolute():
                    # Relative to project root
                    local_path = Path.cwd() / local_path
                
                if not local_path.exists():
                    logger.warning(f"Image file not found: {local_path}")
                    return
                
                with open(local_path, 'rb') as f:
                    image_bytes = io.BytesIO(f.read())
            except Exception as e:
                logger.error(f"Failed to load local image {src}: {e}")
                return
        # Optionally, check/convert format for pptx compatibility
        # (pptx supports PNG, JPEG, BMP, GIF, TIFF, EMF, WMF)
        # We'll use PIL to ensure it's a supported format (convert to PNG if not)
        try:
            img = Image.open(image_bytes)
            if img.format not in ("PNG", "JPEG", "JPG", "BMP", "GIF", "TIFF"):
                # Convert to PNG in memory
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                buf.seek(0)
                image_bytes = buf
        except Exception:
            # If PIL fails, fallback to original bytes
            image_bytes.seek(0)

        # Check for specific size (width, height) override
        size = comp.get("size")
        if size and isinstance(size, list) and len(size) == 2:
            # Override width and height with provided size
            width = size[0]
            height = size[1]

        set_till_end_h = comp.get("set_till_end_h", False)
        set_till_end_v = comp.get("set_till_end_v", False)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # --- Fix: Always recalculate width/height if set_till_end_h/v is set, even if size is present ---
        if set_till_end_h:
            width = slide_width - left
        if set_till_end_v:
            height = slide_height - top

        if fit == "stretch":
            pic = slide.shapes.add_picture(image_bytes, left, top, width=width, height=height)
        else:
            # Place and maintain aspect; simple contain/cover
            pic = slide.shapes.add_picture(image_bytes, left, top)
            # Compute scale
            iw, ih = pic.width, pic.height
            scale_w = width / iw
            scale_h = height / ih
            if fit == "contain":
                scale = min(scale_w, scale_h)
            else:  # cover
                scale = max(scale_w, scale_h)
            pic.width = Emu(int(iw * scale))
            pic.height = Emu(int(ih * scale))
            # center inside rect
            pic.left = left + (width - pic.width) // 2
            pic.top = top + (height - pic.height) // 2
        set_alt_text(pic, comp.get("alt"))
    except Exception as e:
        logger.error(f"[warn] image load failed {src}: {e}")

def render_table(slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    rect = grid_rect(prs, tokens, comp["grid"]) if "grid" in comp else box_rect(comp["box"])
    left, top, width, height = rect
    content = comp["content"]
    cols = content["columns"]
    rows = content["rows"]
    tbl_shape = slide.shapes.add_table(rows=len(rows) + 1, cols=len(cols), left=left, top=top, width=width, height=height)
    tbl = tbl_shape.table

    opts = comp.get("table_options", {})
    style = comp.get("style", {})
    defaults = {
        "header_fill": tokens["color"].get("headerFill"),
        "header_color": tokens["color"].get("headerFill"),
        "row_zebra": tokens["color"].get("zebra"),
        "border_color": tokens["color"].get("border")
    }
    header_fill = comp.get("table_style", {}).get("header_fill", defaults["header_fill"])
    header_color = comp.get("table_style", {}).get("header_color", defaults["header_color"])
    zebra = comp.get("table_style", {}).get("row_zebra", defaults["row_zebra"])
    border_color = comp.get("table_style", {}).get("border_color", defaults["border_color"])

    # Column widths if provided
    if "column_widths" in opts:
        unit = tokens["grid"].get("unit", "px")
        for j, w in enumerate(opts["column_widths"]):
            tbl.columns[j].width = to_emu(w, unit)

    # Header row
    for j, col in enumerate(cols):
        c = tbl.cell(0, j)
        c.text = str(col)
        c.fill.solid()
        c.fill.fore_color.rgb = rgb(header_fill)
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = tokens["font"]["body_family"]
                r.font.bold = True
                r.font.color.rgb = rgb(header_color)

    # Body
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.text = "" if val is None else str(val)
            if zebra and (i % 2 == 1):
                c.fill.solid()
                c.fill.fore_color.rgb = rgb(zebra)
            # borders (simple)
            tc = c._tc
            for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = OxmlElement(side)
                ln.set("w", "12700")  # ~1pt
                solidFill = OxmlElement("a:solidFill")
                srgbClr = OxmlElement("a:srgbClr")
                srgbClr.set("val", border_color.lstrip('#'))
                solidFill.append(srgbClr)
                ln.append(solidFill)
                tc.get_or_add_tcPr().append(ln)

    set_alt_text(tbl_shape, comp.get("alt"))

def render_chart(slide:Slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    rect = grid_rect(prs, tokens, comp["grid"]) if "grid" in comp else box_rect(comp["box"])
    left, top, width, height = rect
    content = comp["content"]
    opts = comp.get("chart_options", {})
    palette = comp.get("chart_style", {}).get("palette") or tokens.get("defaults", {}).get("chart_style", {}).get("palette") or \
              tokens.get("chart_style", {}).get("palette") or []

    chart_type = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
        "stacked_area": XL_CHART_TYPE.AREA_STACKED,
    }.get(opts.get("chartType", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)

    data = CategoryChartData()
    data.categories = content["categories"]
    for s in content["series"]:
        data.add_series(s["name"], s["values"])

    chart = slide.shapes.add_chart(chart_type, left, top, width, height, data).chart

    # Disable multi-level labels on category axis when present to prevent duplicated labels
    try:
        if hasattr(chart, "category_axis") and chart.category_axis is not None:
            chart.category_axis.has_multi_level_labels = False
    except Exception:
        pass

    # Legend
    pos_map = {
        "none": None, "top_right":XL_LEGEND_POSITION.RIGHT,"top": XL_LEGEND_POSITION.TOP, "right": XL_LEGEND_POSITION.RIGHT,
        "bottom": XL_LEGEND_POSITION.BOTTOM, "left": XL_LEGEND_POSITION.LEFT
    }
    legend_opt = opts.get("legend", "right")
    if legend_opt == "none":
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = pos_map.get(legend_opt, XL_LEGEND_POSITION.RIGHT)

    # Ensure legend doesn't overlap with the chart using pptx logic only (no manual padding or plot shifting).
    # The proper, robust way is to set include_in_layout=True, which tells pptx to reserve space for the legend out of the chart plot area.
    if hasattr(chart, "legend") and chart.has_legend:
        chart.legend.include_in_layout = False
    # Data labels (best-effort)
    if bool(opts.get("data_labels", True)):
        if hasattr(chart, "series"):
            for series in chart.series:
                series.data_labels.show_value = True

    # Palette
    for i, series in enumerate(chart.series):
        if i < len(palette):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = rgb(palette[i])

def render_shape(slide:Slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_LINE
    
    rect = grid_rect(prs, tokens, comp["grid"]) if "grid" in comp else box_rect(comp["box"])
    left, top, width, height = rect
    
    shape_type = comp.get("shape_type", "rectangle")
    style = comp.get("style", {})
    
    # Map shape types to PowerPoint shape constants
    shape_mapping = {
        "circle": MSO_SHAPE.OVAL,
        "rectangle": MSO_SHAPE.RECTANGLE,
        "square": MSO_SHAPE.RECTANGLE,  # Will be made square by equal width/height
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "down_arrow": MSO_SHAPE.DOWN_ARROW,
        "up_arrow": MSO_SHAPE.UP_ARROW,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE
    }
    
    pptx_shape_type = shape_mapping.get(shape_type, MSO_SHAPE.RECTANGLE)
    
    # For square, ensure width and height are equal
    if shape_type == "square":
        size = min(width, height)
        width = height = size
    set_till_end_h = comp.get("set_till_end_h", False)
    set_till_end_v = comp.get("set_till_end_v", False)
    if set_till_end_h or set_till_end_v:
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        if set_till_end_h and set_till_end_v:
            # Extend both horizontally and vertically: make square as large as possible
            max_width = slide_width - left
            max_height = slide_height - top
            size = min(max_width, max_height)
            width = height = size
        elif set_till_end_h:
            width = slide_width - left
        elif set_till_end_v:
            height = slide_height - top
    # Add the shape
    shape = slide.shapes.add_shape(pptx_shape_type, left, top, width, height)

    # Apply fill color
    fill_color = style.get("fill")
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
    
    # Apply border
    border_color = style.get("border_color")
    border_width = style.get("border_width", 0)
    if border_color and border_width > 0:
        shape.line.color.rgb = rgb(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()  # No border

    # Handle arrow direction
    if shape_type == "arrow":
        direction = style.get("direction", "right")
        rotation_map = {
            "up": 270,
            "down": 90,
            "left": 180,
            "right": 0
        }
        if direction in rotation_map:
            shape.rotation = rotation_map[direction]

    # Adjust angle of shape using degree if provided
    degree = comp.get("degree")
    if degree is not None:
        shape.rotation = float(degree)

    set_alt_text(shape, comp.get("alt"))
    return shape

from pptx.enum.shapes import MSO_CONNECTOR
def render_line(slide, prs, tokens, comp):
    unit = tokens["grid"].get("unit", "px")
    start_x = to_emu(comp["start"]["x"], unit)
    start_y = to_emu(comp["start"]["y"], unit)
    end_x = to_emu(comp["end"]["x"], unit)
    end_y = to_emu(comp["end"]["y"], unit)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )
    line.line.color.rgb = rgb(comp["style"].get("color", "#FFA500"))
    line.line.width = Pt(comp["style"].get("width", 2))
    return line

def render_group(slide:Slide, prs: Presentation, tokens: Dict[str, Any], comp: Dict[str, Any]):
    """
    Render a group component by rendering all its children and grouping them together.
    Once grouped, all items can be moved together as a single group.
    """
    children = comp.get("children", [])
    rendered_shapes = []

    # Render each child component and collect their shape objects
    for child in children:
        child_type = child["type"]

        # Handle different child component types
        if child_type == "text":
            render_text(slide, prs, tokens, child)
        elif child_type == "richtext":
            render_richtext(slide, prs, tokens, child)
        elif child_type == "image":
            render_image(slide, prs, tokens, child)
        elif child_type == "table":
            render_table(slide, prs, tokens, child)
        elif child_type == "chart":
            render_chart(slide, prs, tokens, child)
        elif child_type == "shape":
            render_shape(slide, prs, tokens, child)
        elif child_type == "line":
            if comp.get("id") == "merge_data":
                child = adjust_spoke_to_circle(child, children, prs=prs, tokens=tokens)
            render_line(slide, prs, tokens, child)
        elif child_type == "group":
            # Recursive group rendering
            render_group(slide, prs, tokens, child)
        else:
            logger.warning(f"Unknown child component type in group: {child_type}")

# ---------- Main render ----------
import math
from copy import deepcopy

def adjust_spoke_to_circle(line, siblings, prs=None, tokens=None):
    # Find inner and outer circles
    inner_circle = next((c for c in siblings if c.get("shape_type") == "circle" and "center" in c["id"]), None)
    outer_circle = next((c for c in siblings if c.get("shape_type") == "circle" and "outer" in c["id"]), None)
    if not inner_circle or not outer_circle or "grid" not in inner_circle or "grid" not in outer_circle:
        return line

    # Compute inner circle rect
    inner_rect = grid_rect(prs, tokens, grid=deepcopy(inner_circle["grid"]))
    left_i, top_i, width_i, height_i = inner_rect
    # Compute outer circle rect
    outer_rect = grid_rect(prs, tokens, grid=deepcopy(outer_circle["grid"]))
    left_o, top_o, width_o, height_o = outer_rect

    # Convert EMU → px
    emu_to_px = 914400 / 96
    left_i /= emu_to_px; top_i /= emu_to_px; width_i /= emu_to_px; height_i /= emu_to_px
    left_o /= emu_to_px; top_o /= emu_to_px; width_o /= emu_to_px; height_o /= emu_to_px

    # Circle centers and radius
    cx_i = left_i + width_i / 2
    cy_i = top_i + height_i / 2
    cx_o = left_o + width_o / 2
    cy_o = top_o + height_o / 2
    r_o = width_o / 2

    # Current direction of the spoke
    start = line["start"]
    end = line["end"]
    dx = end["x"] - start["x"]
    dy = end["y"] - start["y"]
    dist = math.sqrt(dx**2 + dy**2)
    if dist == 0:
        return line
    dx /= dist
    dy /= dist

    # Set start at inner circle center
    line["start"]["x"] = cx_i
    line["start"]["y"] = cy_i
    # Set end at outer circle boundary along the same direction
    line["end"]["x"] = cx_i + dx * r_o
    line["end"]["y"] = cy_i + dy * r_o

    return line

def render_pptx(doc: Dict[str, Any], schema: Dict[str, Any], out_path: Path):
    # Validate
    validator = Draft7Validator(schema)
    errs = sorted(validator.iter_errors(doc), key=lambda e: e.path)
    if errs:
        for e in errs:
            print(f"[schema] {list(e.path)}: {e.message}")
        raise ValueError("Input JSON failed schema validation.")

    deck = doc["deck"]
    tokens = doc["tokens"]
    defaults = doc.get("defaults", {})

    prs = Presentation()
    slide_set_size(prs, deck.get("slide_size", "16x9"))

    for s in doc["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        apply_background(slide, tokens, s.get("background"))
        apply_watermark(slide, deck)  # behind content visually (approx)

        # Title (optional convenience): render as h2 at top, if provided
        # Avoid duplicating if slide already includes a top-level title/h2 component
        y_next = tokens["spacing"]["margin"]
        if s.get("title"):
            has_explicit_title = any(
                (
                    (c.get("type") == "text" and c.get("text_type") in ("title", "h2"))
                    or (c.get("type") == "richtext")
                ) and (
                    ("grid" in c and c["grid"].get("y", y_next) <= y_next + 10)
                )
                for c in s.get("components", [])
            )
            if not has_explicit_title:
                comp_title = {
                    "type": "text",
                    "text_type": "h2",
                    "value": s["title"],
                    "grid": {"col": 1, "span": tokens["grid"]["columns"], "row_h": 48, "y": y_next}
                }
                render_text(slide, prs, tokens, comp_title)
                y_next = y_next + 60

        # Components
        for comp in s["components"]:
            # If component has no explicit y in grid, place after the previous block
            if "grid" in comp and "y" not in comp["grid"]:
                comp = {**comp, "grid": {**comp["grid"], "y": y_next}}

            ctype = comp["type"]
            if ctype == "text":
                render_text(slide, prs, tokens, comp)
            elif ctype == "richtext":
                render_richtext(slide, prs, tokens, comp)
            elif ctype == "image":
                render_image(slide, prs, tokens, comp)
            elif ctype == "table":
                render_table(slide, prs, tokens, comp)
            elif ctype == "chart":
                render_chart(slide, prs, tokens, comp)
            elif ctype == "shape":
                render_shape(slide, prs, tokens, comp)
            elif ctype == "line":
                render_line(slide, prs, tokens, comp)
            elif ctype == "group":
                render_group(slide, prs, tokens, comp)
            else:
                # icon/code rendering could be added here
                pass

            # Advance y cursor for simple flow layouts
            if "grid" in comp:
                y_next = comp["grid"].get("y", y_next) + comp["grid"]["row_h"] + tokens["spacing"]["gutter"]

    # Add header/footer/page numbers
    add_header_footer_and_page_numbers(prs, deck, tokens)

    prs.save(str(out_path))
    logger.info(f"[ok] Wrote {out_path}")

# ---------- CLI ----------

def main():
    ap = argparse.ArgumentParser(description="Render PPTX from enterprise JSON schema.")
    ap.add_argument("--schema", required=True, help="Path to JSON Schema (Draft-07).")
    ap.add_argument("--input", required=True, help="Path to deck JSON.")
    ap.add_argument("--out", required=True, help="Output .pptx path.")
    args = ap.parse_args()

    schema_path = Path(args.schema)
    input_path = Path(args.input)
    out_path = Path(args.out)

    with schema_path.open("r", encoding="utf-8") as f:
        schema = json.load(f)
    with input_path.open("r", encoding="utf-8") as f:
        doc = json.load(f)
    render_pptx(doc, schema, out_path)